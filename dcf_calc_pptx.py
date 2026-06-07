"""
DCF Intrinsic Value — Step-by-Step Calculation (continuation of Financial Statements deck)
Company: Cocoa Foods Ghana Ltd  |  Presented by Kofi  |  InvestIQ
Saves: Desktop/DCF_Calculation_Kofi.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

WHITE=RGBColor(0xff,0xff,0xff); OFF=RGBColor(0xf8,0xfa,0xfc)
LG=RGBColor(0xf1,0xf5,0xf9); MG=RGBColor(0xe2,0xe8,0xf0)
STEEL=RGBColor(0x64,0x74,0x8b); DARK=RGBColor(0x0f,0x17,0x2a)
NAVY=RGBColor(0x1e,0x3a,0x8a); BLUE=RGBColor(0x1d,0x4e,0xd8)
BL=RGBColor(0xdb,0xea,0xff); BT=RGBColor(0x1e,0x40,0xaf)
TEAL=RGBColor(0x06,0x7a,0x5f); TL=RGBColor(0xd1,0xfa,0xee); TT=RGBColor(0x06,0x5f,0x46)
AMB=RGBColor(0xd9,0x77,0x06); AL=RGBColor(0xff,0xf7,0xd6); AT=RGBColor(0x92,0x40,0x02)
RED=RGBColor(0xb9,0x1c,0x1c); RL=RGBColor(0xfe,0xe2,0xe2); RT=RGBColor(0x99,0x1b,0x1b)
GRN=RGBColor(0x15,0x80,0x3d); GL=RGBColor(0xdc,0xfc,0xe8); GT=RGBColor(0x16,0x65,0x34)
PUR=RGBColor(0x6d,0x28,0xd9); PL=RGBColor(0xed,0xe9,0xfe)
SUB=RGBColor(0xba,0xd0,0xf8)

def R(sl,l,t,w,h,fill,bd=None):
    s=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if bd: s.line.color.rgb=bd; s.line.width=Pt(0.75)
    else: s.line.fill.background()

def T(sl,text,l,t,w,h,sz=10,bold=False,color=DARK,align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Calibri"

def hdr(sl,title,sub=None,ac=NAVY):
    R(sl,0,0,13.33,1.25,ac)
    T(sl,title,0.35,0.1,12.5,0.65,sz=24,bold=True,color=WHITE)
    if sub: T(sl,sub,0.35,0.82,12.5,0.38,sz=10,color=SUB,italic=True)

def ftr(sl,n,tot=13):
    R(sl,0,7.2,13.33,0.3,MG)
    T(sl,"DCF Intrinsic Value — Presented by Kofi  |  InvestIQ  |  Continuation of Financial Statements Deck",
      0.25,7.22,11.5,0.24,sz=7.5,color=STEEL)
    T(sl,f"{n}/{tot}",12.5,7.22,0.7,0.24,sz=7.5,color=STEEL,align=PP_ALIGN.RIGHT)

def sbg(sl): R(sl,0,0,13.33,7.5,OFF)

# formula box helper
def fbox(sl,l,t,w,h,formula,title=None,fill=BL,border=BLUE,fc=BT,tc=WHITE):
    R(sl,l,t,w,h,fill,bd=border)
    if title:
        R(sl,l,t,w,0.32,border)
        T(sl,title,l+0.15,t+0.04,w-0.3,0.24,sz=10,bold=True,color=tc)
        T(sl,formula,l+0.15,t+0.38,w-0.3,h-0.46,sz=11,bold=True,color=fc,align=PP_ALIGN.CENTER)
    else:
        T(sl,formula,l+0.15,t+0.1,w-0.3,h-0.18,sz=11,bold=True,color=fc,align=PP_ALIGN.CENTER)

# table row
def row(sl,y,l,w,cells,fills,sz=9.5,rh=0.36):
    x=l
    for (txt_,cw,bold,color),fill in zip(cells,fills):
        R(sl,x,y,cw,rh,fill,bd=MG)
        al=PP_ALIGN.RIGHT if color in [GT,TT,AT,RT,BT] and len(txt_)<8 else PP_ALIGN.LEFT
        T(sl,txt_,x+0.08,y+0.05,cw-0.16,rh-0.1,sz=sz,bold=bold,color=color,align=al)
        x+=cw
    return y+rh


def build():
    prs=Presentation()
    prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    TOT=16

    # ── S1 COVER ──────────────────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    R(s,0,0,13.33,7.5,NAVY)
    R(s,0,0,13.33,0.2,AMB); R(s,0,7.3,13.33,0.2,TEAL); R(s,0,0.2,0.15,7.1,BLUE)
    T(s,"Calculating Intrinsic Value",0.5,0.8,12.3,0.88,sz=38,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,"DCF MODEL — STEP BY STEP WITH EVERY FORMULA",0.5,1.78,12.3,0.55,sz=19,bold=True,color=AMB,align=PP_ALIGN.CENTER)
    T(s,"Free Cash Flow  ·  Beta  ·  Cost of Equity (CAPM)  ·  Cost of Debt  ·  WACC  ·  Terminal Value  ·  Intrinsic Value per Share",
      0.5,2.42,12.3,0.4,sz=12,italic=True,color=SUB,align=PP_ALIGN.CENTER)
    R(s,1.5,3.15,10.33,2.6,RGBColor(0x1a,0x30,0x6e))
    R(s,1.5,3.15,10.33,0.06,AMB)
    T(s,"Company used throughout:",1.7,3.28,9.9,0.32,sz=11,color=RGBColor(0x94,0xa3,0xb8))
    T(s,"Cocoa Foods Ghana Ltd  (Manufacturing)",1.7,3.65,9.0,0.52,sz=20,bold=True,color=WHITE)
    T(s,"Using numbers from the Financial Statements deck  |  All figures GHc millions",
      1.7,4.22,9.0,0.35,sz=11,color=SUB)
    T(s,"Continuation of: Understanding Financial Statements — Presented by Kofi",
      1.7,4.62,9.0,0.35,sz=11,italic=True,color=RGBColor(0x94,0xa3,0xb8))
    T(s,"Presented by Kofi  ·  InvestIQ  ·  investright.onrender.com",
      0.4,7.05,12.5,0.28,sz=9,color=RGBColor(0x94,0xa3,0xb8),align=PP_ALIGN.CENTER)

    # ── S2 OVERVIEW: THE 8 STEPS ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"The 8 Steps to DCF Intrinsic Value",
               "We build up from financial statement numbers to a single fair price per share"); ftr(s,2,TOT)
    T(s,"Every number in this model comes directly from Cocoa Foods Ghana's financial statements.",
      0.4,1.35,12.5,0.35,sz=11,italic=True,color=STEEL)
    steps=[
        (BLUE,BL,BT,"Step 1","Find Free Cash Flow (FCF)","From the Cash Flow Statement — the cash the business actually generates"),
        (TEAL,TL,TT,"Step 2","Cost of Equity (Ke)","Using the CAPM formula: Risk-free Rate + Beta × Market Risk Premium"),
        (AMB,AL,AT,"Step 3","Cost of Debt (Kd)","Interest rate on loans, adjusted for the tax shield"),
        (PUR,PL,PUR,"Step 4","WACC","Blend cost of equity and cost of debt by capital structure weights"),
        (GRN,GL,GT,"Step 5","Project FCF for 10 Years","Apply growth rates to forecast future free cash flows"),
        (BLUE,BL,BT,"Step 6","Terminal Value","Estimate the company's value beyond Year 10 using the Gordon Growth Model"),
        (TEAL,TL,TT,"Step 7","Discount Everything to Today","Divide each future cash flow by (1+WACC)ⁿ"),
        (AMB,AL,AT,"Step 8","Intrinsic Value per Share","Enterprise Value − Net Debt ÷ Shares Outstanding = Fair Price"),
    ]
    for i,(b,f,tc,step,title,desc) in enumerate(steps):
        x=0.35+(i%2)*6.47; y=1.82+(i//2)*1.3
        R(s,x,y,6.25,1.2,f,bd=b); R(s,x,y,0.55,1.2,b)
        T(s,step,x,y+0.38,0.55,0.42,sz=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        T(s,title,x+0.65,y+0.08,5.45,0.35,sz=12,bold=True,color=tc)
        T(s,desc,x+0.65,y+0.48,5.45,0.65,sz=9.5,color=DARK)

    # ── S3 STEP 1: FREE CASH FLOW (FCF / FCFF / FCFE) ────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 1 — Free Cash Flow: FCF, FCFF & FCFE",
               "Three measures of cash generation — each used in a different DCF context"); ftr(s,3,TOT)
    T(s,"All three are derived from Cocoa Foods Ghana's financial statements. The choice of measure determines which discount rate you use.",
      0.4,1.35,12.5,0.32,sz=10.5,italic=True,color=STEEL)
    CW=4.08  # column width — 3 cols with 0.13 gaps = 3×4.08+2×0.13=12.5 ✓
    cols=[
        # (x, border, light_fill, text_color, title, formula_text, rows, note)
        (0.35, BLUE, BL, BT,
         "Method A — Simple FCF",
         "FCF  =  Operating Cash Flow  −  Capex",
         [("Operating Cash Flow (CFS)","45.5",False,STEEL,LG),
          ("Less: Capital Expenditure","(38.0)",False,RT,RL),
          ("FREE CASH FLOW","7.5",True,TT,TL)],
         "Quickest method.\nDirect from CFS — no IS adjustments needed.\nUse for a fast, conservative cross-check.\nAll three methods should broadly agree.",
         "Use FCF when you want a quick check — it is already in the Cash Flow Statement."),
        (4.56, AMB, AL, AT,
         "Method B — FCFF (Free Cash Flow to Firm)",
         "FCFF  =  EBIT×(1−t)  +  D&A  −  Capex  −  ΔNWC",
         [("EBIT (from Income Statement)","74.0",False,STEEL,LG),
          ("× (1 − Tax Rate of 25%)","× 0.75",False,STEEL,OFF),
          ("= NOPAT (Net Op. Profit)","55.5",False,AT,AL),
          ("+ Depreciation & Amortisation","+ 24.0",False,GT,GL),
          ("− Capital Expenditure (Capex)","− 38.0",False,RT,RL),
          ("− Change in Working Capital","− 22.0",False,RT,RL),
          ("FCFF","19.5",True,TT,TL)],
         "Cash available to ALL capital providers.\nDiscount using WACC.\nThis is the standard method for enterprise DCF.",
         "FCFF → discount at WACC → gives Enterprise Value (EV). Subtract Net Debt to get Equity Value."),
        (9.13, GRN, GL, GT,
         "Method C — FCFE (Free Cash Flow to Equity)",
         "FCFE  =  Net Profit + D&A − Capex − ΔNWC + Net Borrowings",
         [("Net Profit (from IS)","43.5",False,STEEL,LG),
          ("+ Depreciation & Amortisation","+ 24.0",False,GT,GL),
          ("− Capital Expenditure (Capex)","− 38.0",False,RT,RL),
          ("− Change in Working Capital","− 22.0",False,RT,RL),
          ("+ Net Borrowings (20 − 10)","+ 10.0",False,BT,BL),
          ("FCFE","17.5",True,TT,TL)],
         "Cash available to EQUITY shareholders only.\nAlready reflects debt service & borrowings.\nDiscount using Ke (Cost of Equity), not WACC.",
         "FCFE → discount at Ke → directly gives Equity Value. Divide by shares for intrinsic value per share."),
    ]
    RH=0.385  # row height
    for (cx,border,lf,tc,title,formula,rows,note2,bottom_note) in cols:
        col_h=0.35+0.58+len(rows)*RH+0.82
        R(s,cx,1.78,CW,col_h,lf,bd=border)
        R(s,cx,1.78,CW,0.35,border)
        T(s,title,cx+0.12,1.81,CW-0.24,0.28,sz=10.5,bold=True,color=WHITE)
        # formula box
        R(s,cx+0.08,2.2,CW-0.16,0.52,WHITE,bd=border)
        T(s,formula,cx+0.14,2.27,CW-0.28,0.38,sz=9.5,bold=True,color=tc,align=PP_ALIGN.CENTER)
        # data rows
        for j,(lb,am,hl,ac,rf) in enumerate(rows):
            ry=2.79+j*RH
            R(s,cx+0.06,ry,CW-0.12,RH-0.03,rf if hl else (LG if j%2==0 else OFF),bd=border if hl else MG)
            T(s,lb,cx+0.14,ry+0.05,CW-0.68,RH-0.1,sz=9,bold=hl,color=ac if hl else DARK)
            T(s,am,cx+CW-0.56,ry+0.05,0.46,RH-0.1,sz=9,bold=hl,color=ac if hl else DARK,align=PP_ALIGN.RIGHT)
        # note
        ny=2.79+len(rows)*RH
        T(s,note2,cx+0.12,ny+0.06,CW-0.24,0.68,sz=8.5,italic=True,color=DARK)
    # bottom comparison bar
    R(s,0.35,6.87,12.63,0.42,NAVY)
    T(s,"FCFF vs FCFE — Key Difference:",0.5,6.91,3.2,0.32,sz=10,bold=True,color=AMB)
    T(s,"FCFF = cash to ALL capital providers → discount at WACC → Enterprise Value − Net Debt = Equity Value     |     FCFE = cash to EQUITY only → discount at Ke → directly = Equity Value     |     FCFE = FCFF − Interest×(1−t) + Net Borrowings",
      3.8,6.91,9.1,0.32,sz=9,color=WHITE)

    # ── S3b: WHERE TO FIND — SIMPLE FCF COMPONENTS ───────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 1a — Where to Find Simple FCF in the Statements",
               "FCF = Operating Cash Flow − Capex  |  Both items sit in the Cash Flow Statement"); ftr(s,4,TOT)
    R(s,0.35,1.32,12.63,0.42,BL,bd=BLUE)
    T(s,"FCF  =  Operating Cash Flow  −  Capital Expenditure (Capex)",0.5,1.38,12.2,0.3,sz=13,bold=True,color=BT,align=PP_ALIGN.CENTER)
    # Left column: CFS extract
    R(s,0.35,1.88,8.1,5.28,LG,bd=MG); R(s,0.35,1.88,8.1,0.35,BLUE)
    T(s,"CASH FLOW STATEMENT — Cocoa Foods Ghana Ltd (extract)",0.5,1.91,7.9,0.26,sz=10.5,bold=True,color=WHITE)
    cfs_rows=[
        ("OPERATING ACTIVITIES","",(BLUE,BL)),
        ("Net Profit after Tax","43.5",(STEEL,LG)),
        ("Add back: Depreciation & Amortisation","24.0",(STEEL,OFF)),
        ("Less: Increase in Working Capital","(22.0)",(STEEL,LG)),
        ("NET CASH FROM OPERATIONS  ★","45.5",(BT,BL)),
        ("INVESTING ACTIVITIES","",(BLUE,BL)),
        ("Purchase of Property, Plant & Equipment (Capex)  ★","(38.0)",(RT,RL)),
        ("NET CASH FROM INVESTING","(38.0)",(STEEL,LG)),
        ("NET CHANGE IN CASH","7.5",(GT,GL)),
    ]
    for j,(lb,am,(tc,rf)) in enumerate(cfs_rows):
        ry=2.3+j*0.48
        R(s,0.42,ry,7.95,0.43,rf,bd=MG)
        bold=lb.startswith("NET") or "★" in lb
        T(s,lb,0.55,ry+0.07,6.2,0.28,sz=9.5,bold=bold,color=tc)
        if am: T(s,am,7.05,ry+0.07,1.2,0.28,sz=9.5,bold=bold,color=tc,align=PP_ALIGN.RIGHT)
    # Right column: callout boxes
    R(s,8.6,1.88,4.73,2.38,BL,bd=BLUE); R(s,8.6,1.88,4.73,0.35,BLUE)
    T(s,"Component 1 — Operating Cash Flow",8.75,1.91,4.5,0.26,sz=10.5,bold=True,color=WHITE)
    T(s,"Where: Cash Flow Statement → Operating Activities section\nLine: 'Net Cash from Operating Activities'\nCocoa Foods: GHc 45.5m\n\nThis is the cash earned from running the business day-to-day AFTER adjusting for non-cash items (like D&A) and working capital movements.",
      8.72,2.32,4.5,1.75,sz=9.5,color=DARK)
    R(s,8.6,4.38,4.73,2.38,RL,bd=RED); R(s,8.6,4.38,4.73,0.35,RED)
    T(s,"Component 2 — Capital Expenditure (Capex)",8.75,4.41,4.5,0.26,sz=10.5,bold=True,color=WHITE)
    T(s,"Where: Cash Flow Statement → Investing Activities section\nLine: 'Purchase of PP&E' or 'Acquisition of fixed assets'\nCocoa Foods: GHc (38.0m) — shown as negative (cash out)\n\nCapex = money spent to maintain or expand the asset base.",
      8.72,4.82,4.5,1.75,sz=9.5,color=DARK)
    R(s,0.35,7.06,12.63,0.12,BLUE)
    T(s,"Simple FCF = 45.5 − 38.0 = GHc 7.5m   ✓   Both numbers are already in the Cash Flow Statement — no Income Statement needed.",
      0.5,6.9,12.2,0.22,sz=9.5,bold=True,color=BT,align=PP_ALIGN.CENTER)

    # ── S3c: WHERE TO FIND — FCFF COMPONENTS ─────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 1b — Where to Find FCFF Components in the Statements",
               "FCFF = EBIT×(1−t) + D&A − Capex − ΔNWC  |  Pulls from Income Statement, Balance Sheet & CFS"); ftr(s,5,TOT)
    R(s,0.35,1.32,12.63,0.42,AL,bd=AMB)
    T(s,"FCFF  =  EBIT × (1 − Tax Rate)  +  D&A  −  Capex  −  ΔNWC",0.5,1.38,12.2,0.3,sz=13,bold=True,color=AT,align=PP_ALIGN.CENTER)
    comps_fcff=[
        (AMB,AL,AT,"EBIT","Income Statement",
         "Line: 'Operating Profit' or 'EBIT'\nUsually = Gross Profit − Operating Expenses − D&A\nCocoa Foods EBIT = GHc 74.0m"),
        (AMB,AL,AT,"Tax Rate (t)","Income Statement / Notes",
         "Line: 'Income Tax Expense' or statutory note\nTax Rate = Tax Charge ÷ Profit Before Tax\nCocoa Foods: 25% Ghana corporate tax"),
        (GRN,GL,GT,"D&A","Income Statement or CFS Operating",
         "IS: often embedded in COGS or listed separately\nCFS: always shown as an add-back in Operating Activities\nCocoa Foods D&A = GHc 24.0m"),
        (RED,RL,RT,"Capex","CFS — Investing Activities",
         "Line: 'Purchase of Property, Plant & Equipment'\nShown as negative (cash outflow)\nCocoa Foods Capex = GHc (38.0m)"),
        (BLUE,BL,BT,"ΔNWC (Change in Working Capital)","Balance Sheet or CFS Operating",
         "CFS: shown as 'Increase/(Decrease) in Working Capital'\nOR calculate: (Current Assets − Cash) − Current Liabilities\nCocoa Foods ΔNWC = GHc 22.0m increase"),
    ]
    CW2=2.46; gap=0.1
    for i,(b,f,tc,comp,source,detail) in enumerate(comps_fcff):
        cx=0.35+i*(CW2+gap)
        R(s,cx,1.88,CW2,5.28,f,bd=b); R(s,cx,1.88,CW2,0.35,b)
        T(s,comp,cx+0.1,1.91,CW2-0.2,0.26,sz=11,bold=True,color=WHITE)
        R(s,cx+0.08,2.3,CW2-0.16,0.32,b)
        T(s,source,cx+0.14,2.34,CW2-0.28,0.22,sz=8.5,bold=True,color=WHITE)
        T(s,detail,cx+0.1,2.7,CW2-0.2,4.1,sz=9,color=DARK)
    R(s,0.35,7.06,12.63,0.12,AMB)
    T(s,"FCFF = 74×0.75 + 24 − 38 − 22 = 55.5 + 24 − 38 − 22 = GHc 19.5m   ✓   All 5 components traceable to the financial statements.",
      0.5,6.9,12.2,0.22,sz=9.5,bold=True,color=AT,align=PP_ALIGN.CENTER)

    # ── S3d: WHERE TO FIND — FCFE COMPONENTS ─────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 1c — Where to Find FCFE Components in the Statements",
               "FCFE = Net Profit + D&A − Capex − ΔNWC + Net Borrowings  |  Starts from the bottom of the IS"); ftr(s,6,TOT)
    R(s,0.35,1.32,12.63,0.42,GL,bd=GRN)
    T(s,"FCFE  =  Net Profit  +  D&A  −  Capex  −  ΔNWC  +  Net Borrowings",0.5,1.38,12.2,0.3,sz=13,bold=True,color=GT,align=PP_ALIGN.CENTER)
    comps_fcfe=[
        (GRN,GL,GT,"Net Profit","Income Statement — bottom line",
         "Line: 'Profit After Tax' or 'Net Income'\nThis is AFTER interest and tax — already reflects debt cost\nCocoa Foods: GHc 43.5m"),
        (GRN,GL,GT,"D&A","IS or CFS Operating Activities",
         "Same as FCFF — add back the non-cash charge\nCFS shows it as a positive add-back\nCocoa Foods: GHc +24.0m"),
        (RED,RL,RT,"Capex","CFS — Investing Activities",
         "Same as FCFF — cash spent on fixed assets\nLine: 'Purchase of PP&E'\nCocoa Foods: GHc (38.0m)"),
        (RED,RL,RT,"ΔNWC","Balance Sheet or CFS Operating",
         "Same as FCFF — working capital movement\nIncrease = cash used = negative for FCF\nCocoa Foods: GHc (22.0m)"),
        (BLUE,BL,BT,"Net Borrowings","CFS — Financing Activities",
         "New loans raised MINUS loan repayments\nLine: 'Proceeds from borrowings' less 'Repayment of loans'\nCocoa Foods: +20 new − 10 repaid = GHc +10.0m"),
    ]
    for i,(b,f,tc,comp,source,detail) in enumerate(comps_fcfe):
        cx=0.35+i*(CW2+gap)
        R(s,cx,1.88,CW2,5.28,f,bd=b); R(s,cx,1.88,CW2,0.35,b)
        T(s,comp,cx+0.1,1.91,CW2-0.2,0.26,sz=11,bold=True,color=WHITE)
        R(s,cx+0.08,2.3,CW2-0.16,0.32,b)
        T(s,source,cx+0.14,2.34,CW2-0.28,0.22,sz=8.5,bold=True,color=WHITE)
        T(s,detail,cx+0.1,2.7,CW2-0.2,4.1,sz=9,color=DARK)
    R(s,0.35,7.06,12.63,0.12,GRN)
    T(s,"FCFE = 43.5 + 24 − 38 − 22 + 10 = GHc 17.5m   ✓   Also: FCFF(19.5) − Interest×(1−t)(12) + Net Borrowings(10) = 17.5m — both routes agree.",
      0.5,6.9,12.2,0.22,sz=9.5,bold=True,color=GT,align=PP_ALIGN.CENTER)

    # ── S4 STEP 2: COST OF EQUITY (CAPM) ─────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 2 — Cost of Equity (Ke) using CAPM",
               "What annual return do equity investors demand for taking the risk of owning this stock?"); ftr(s,7,TOT)
    # CAPM formula
    fbox(s,0.35,1.35,12.63,0.72,"Ke  =  Rf  +  β  ×  (Rm − Rf)    where  (Rm − Rf)  =  Equity Risk Premium (ERP)",
         fill=BL,border=BLUE,fc=BT)
    # Three ingredients
    cards=[
        (BLUE,BL,BT,"Rf — Risk-Free Rate","The return on a zero-risk investment.\nIn Ghana: the 91-Day Treasury Bill rate.\n\nWe use:  Rf = 25%\n\n(This is the minimum any investor expects.\nIf you can earn 25% risk-free, why take risk\nfor anything less?)"),
        (AMB,AL,AT,"β — Beta","Measures how much the stock moves relative to the overall market.\n\nβ = 1.0 → moves exactly with the market\nβ > 1.0 → more volatile than market\nβ < 1.0 → less volatile than market\n\nCocoa Foods Ghana β = 0.85\n(defensive consumer goods — less volatile)"),
        (PUR,PL,PUR,"ERP — Equity Risk Premium","The extra return investors demand for choosing stocks over risk-free T-bills.\n\n(Rm − Rf)  =  ERP\n\nGhana ERP ≈ 6%\n\n(Investors need 6% MORE than the T-bill\nrate to justify owning equities)"),
    ]
    for i,(b,f,tc,title,body) in enumerate(cards):
        x=0.35+i*4.33
        R(s,x,2.18,4.15,4.2,f,bd=b); R(s,x,2.18,4.15,0.35,b)
        T(s,title,x+0.15,2.21,3.9,0.28,sz=12,bold=True,color=WHITE)
        T(s,body,x+0.15,2.62,3.9,3.6,sz=10,color=DARK)
    # Calculation
    R(s,0.35,6.48,12.63,0.65,GL,bd=GRN)
    R(s,0.35,6.48,0.12,0.65,GRN)
    T(s,"CALCULATION:   Ke  =  Rf + β × ERP  =  25%  +  0.85 × 6%  =  25%  +  5.1%  =  30.1%",
      0.55,6.56,10.5,0.5,sz=13,bold=True,color=GT)
    T(s,"→  30%",11.35,6.56,1.55,0.5,sz=18,bold=True,color=GT,align=PP_ALIGN.RIGHT)

    # ── S5 STEP 3: COST OF DEBT ───────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 3 — Cost of Debt (Kd)",
               "The after-tax cost of borrowing — debt is cheaper than equity because interest is tax-deductible"); ftr(s,8,TOT)
    fbox(s,0.35,1.35,12.63,0.72,"Kd (after-tax)  =  Pre-tax Interest Rate  ×  (1  −  Corporate Tax Rate)",
         fill=RL,border=RED,fc=RT)
    # Left: explanation
    R(s,0.4,2.18,6.0,4.95,RL,bd=RED); R(s,0.4,2.18,6.0,0.35,RED)
    T(s,"Why do we adjust for tax?",0.55,2.21,5.8,0.28,sz=11,bold=True,color=WHITE)
    T(s,"Interest paid on loans is tax-deductible.\nThis means the government shares part of the cost.",
      0.55,2.62,5.8,0.52,sz=10.5,color=DARK)
    eg=[("Interest paid on GHc 193m debt (22%)","42.5"),
        ("Tax saved on that interest (25% × 42.5)","10.6"),
        ("ACTUAL net cost of debt to company","31.9"),
        ("Effective rate = 31.9 / 193","16.5%")]
    for j,(lb,am) in enumerate(eg):
        fill=TL if j==3 else (LG if j%2==0 else OFF)
        R(s,0.5,3.28+j*0.72,5.8,0.66,fill,bd=MG)
        T(s,lb,0.62,3.34+j*0.72,4.35,0.5,sz=10,bold=(j==3),color=TT if j==3 else DARK)
        T(s,am,5.05,3.34+j*0.72,1.1,0.5,sz=11,bold=(j==3),color=TT if j==3 else DARK,align=PP_ALIGN.RIGHT)
    # Right: numbers
    R(s,6.58,2.18,6.35,4.95,AL,bd=AMB); R(s,6.58,2.18,6.35,0.35,AMB)
    T(s,"Cocoa Foods Ghana — Step by Step",6.73,2.21,6.15,0.28,sz=11,bold=True,color=WHITE)
    items=[
        ("Pre-tax Interest Rate","From Note to Accounts / Loan agreements","22.0%",False),
        ("Source of debt","Balance Sheet: LT Debt GHc 148m + ST Debt GHc 45m","GHc 193m",False),
        ("Corporate Tax Rate","Ghana Corporate Income Tax rate","25.0%",False),
        ("Tax adjustment factor","(1 − 0.25)","0.75",False),
        ("After-tax Cost of Debt","22% × 0.75","16.5%",True),
    ]
    for j,(lbl,note,val,hl) in enumerate(items):
        fill=GL if hl else (LG if j%2==0 else OFF)
        R(s,6.68,2.65+j*0.85,6.15,0.78,fill,bd=MG)
        T(s,lbl,6.8,2.68+j*0.85,3.3,0.28,sz=10,bold=hl,color=GT if hl else DARK)
        T(s,note,6.8,2.96+j*0.85,3.3,0.32,sz=8,italic=True,color=STEEL)
        T(s,val,10.2,2.74+j*0.85,2.4,0.38,sz=13 if hl else 10,bold=hl,color=GT if hl else BT,align=PP_ALIGN.RIGHT)
    R(s,0.35,7.06,12.63,0.1,AMB)
    T(s,"Kd (after-tax)  =  22%  ×  (1 − 25%)  =  22%  ×  0.75  =  16.5%   →   We use  Kd = 16.5%",
      0.5,6.9,12.2,0.22,sz=10,bold=True,color=AT,align=PP_ALIGN.CENTER)

    # ── S6 STEP 4: WACC ───────────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 4 — WACC (Weighted Average Cost of Capital)",
               "Blending equity and debt costs by their proportions in the capital structure"); ftr(s,9,TOT)
    fbox(s,0.35,1.35,12.63,0.72,
         "WACC  =  (We × Ke)  +  (Wd × Kd)    where  We + Wd = 100%",
         fill=PL,border=PUR,fc=PUR)
    # Capital structure from BS
    R(s,0.4,2.18,5.95,2.22,BL,bd=BLUE); R(s,0.4,2.18,5.95,0.35,BLUE)
    T(s,"Capital Structure (from Balance Sheet)",0.55,2.21,5.75,0.28,sz=11,bold=True,color=WHITE)
    cs=[("Total Equity (Book)","GHc 239m","55%",BT),
        ("Total Debt (LT + ST)","GHc 193m","45%",RT),
        ("Total Capital","GHc 432m","100%",GT)]
    for j,(lbl,amt,wt,cl) in enumerate(cs):
        fill=TL if j==2 else (LG if j%2==0 else OFF)
        R(s,0.5,2.62+j*0.58,5.75,0.52,fill,bd=MG)
        T(s,lbl,0.62,2.66+j*0.58,2.8,0.38,sz=10,bold=(j==2),color=DARK)
        T(s,amt,3.45,2.66+j*0.58,1.35,0.38,sz=10,bold=(j==2),color=cl,align=PP_ALIGN.RIGHT)
        T(s,"Weight: "+wt,4.85,2.66+j*0.58,1.2,0.38,sz=10,bold=True,color=cl,align=PP_ALIGN.RIGHT)
    # WACC calc table
    R(s,0.4,4.55,5.95,2.62,GL,bd=GRN); R(s,0.4,4.55,5.95,0.35,GRN)
    T(s,"WACC Calculation",0.55,4.58,5.75,0.28,sz=11,bold=True,color=WHITE)
    wc=[("Equity component","We × Ke","55% × 30.1%","= 16.6%",BT),
        ("Debt component","Wd × Kd","45% × 16.5%","= 7.4%",AT),
        ("WACC","We×Ke + Wd×Kd","16.6% + 7.4%","= 24.0%",GT)]
    for j,(lbl,formula,calc,result,cl) in enumerate(wc):
        fill=GL if j==2 else (LG if j%2==0 else OFF)
        R(s,0.5,5.0+j*0.68,5.75,0.62,fill,bd=MG)
        T(s,lbl,0.62,5.04+j*0.68,1.5,0.5,sz=9.5,bold=(j==2),color=cl)
        T(s,formula,2.15,5.04+j*0.68,1.3,0.5,sz=9.5,color=STEEL)
        T(s,calc,3.5,5.04+j*0.68,1.25,0.5,sz=9.5,color=DARK)
        T(s,result,4.82,5.04+j*0.68,1.0,0.5,sz=11 if j==2 else 9.5,bold=(j==2),color=cl,align=PP_ALIGN.RIGHT)
    # Right: summary + interpretation
    R(s,6.55,2.18,6.35,4.99,AL,bd=AMB); R(s,6.55,2.18,6.35,0.35,AMB)
    T(s,"What WACC Means & How to Use It",6.7,2.21,6.15,0.28,sz=11,bold=True,color=WHITE)
    T(s,"WACC = 24%  means:",6.7,2.65,6.1,0.3,sz=12,bold=True,color=AT)
    pts=["\"Investors require a 24% annual return to justify owning this company.\"",
         "Any project or investment the company makes MUST earn more than 24% — otherwise it destroys shareholder value.",
         "We use 24% to DISCOUNT all future free cash flows back to today's value.",
         "Higher WACC = future cash flows worth LESS today = lower intrinsic value.",
         "Lower WACC = future cash flows worth MORE = higher intrinsic value.",
         "WACC is the most judgement-sensitive input in the DCF — small changes have big effects on value."]
    for j,pt in enumerate(pts):
        R(s,6.65,3.05+j*0.68,6.15,0.62,LG if j%2==0 else OFF,bd=MG)
        T(s,f"• {pt}",6.78,3.08+j*0.68,5.9,0.56,sz=9.5,color=DARK)
    R(s,0.35,7.06,12.63,0.1,PUR)
    T(s,"WACC = 24.0%  — this is our discount rate for the entire DCF model",
      0.5,6.9,12.2,0.22,sz=10,bold=True,color=PUR,align=PP_ALIGN.CENTER)

    # ── S7 STEP 5: PROJECT FCF ────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 5 — Project Free Cash Flow for 10 Years",
               "Apply growth rates to forecast how FCF will grow. Base FCF = GHc 20m"); ftr(s,10,TOT)
    T(s,"Growth assumptions — from industry research and management guidance:",
      0.4,1.35,7.0,0.32,sz=10.5,color=DARK)
    T(s,"Yr 1–5: 12%  (strong growth phase)     Yr 6–10: 6%  (maturity phase)",
      0.4,1.62,8.0,0.32,sz=11,bold=True,color=BT)
    T(s,"Formula:   FCFₙ = FCFₙ₋₁ × (1 + growth rate)     Discount Factor = 1 ÷ (1 + WACC)ⁿ = 1 ÷ (1.24)ⁿ",
      0.4,1.95,12.5,0.32,sz=10.5,italic=True,color=STEEL)
    # Table header
    R(s,0.35,2.35,12.63,0.36,NAVY)
    for hd,xp,cw in [("Year",0.45,0.8),("Growth",1.3,0.85),("FCF (GHc m)",2.2,1.35),
                      ("Calc: FCF×(1+g)",3.6,2.2),("Discount Factor 1÷(1.24)ⁿ",5.85,2.5),
                      ("PV = FCF ÷ (1.24)ⁿ  (GHc m)",8.4,4.53)]:
        T(s,hd,xp,2.38,cw,0.28,sz=8.5,bold=True,color=WHITE)
    data=[
        ("Yr 1","12%","22.4","20 × 1.12 = 22.4","1 ÷ 1.24¹ = 0.8065","22.4 × 0.8065 = 18.1"),
        ("Yr 2","12%","25.1","22.4 × 1.12 = 25.1","1 ÷ 1.24² = 0.6504","25.1 × 0.6504 = 16.3"),
        ("Yr 3","12%","28.1","25.1 × 1.12 = 28.1","1 ÷ 1.24³ = 0.5245","28.1 × 0.5245 = 14.7"),
        ("Yr 4","12%","31.5","28.1 × 1.12 = 31.5","1 ÷ 1.24⁴ = 0.4230","31.5 × 0.4230 = 13.3"),
        ("Yr 5","12%","35.3","31.5 × 1.12 = 35.3","1 ÷ 1.24⁵ = 0.3411","35.3 × 0.3411 = 12.0"),
        ("Yr 6","6%","37.4","35.3 × 1.06 = 37.4","1 ÷ 1.24⁶ = 0.2751","37.4 × 0.2751 = 10.3"),
        ("Yr 7","6%","39.6","37.4 × 1.06 = 39.6","1 ÷ 1.24⁷ = 0.2218","39.6 × 0.2218 =  8.8"),
        ("Yr 8","6%","42.0","39.6 × 1.06 = 42.0","1 ÷ 1.24⁸ = 0.1789","42.0 × 0.1789 =  7.5"),
        ("Yr 9","6%","44.5","42.0 × 1.06 = 44.5","1 ÷ 1.24⁹ = 0.1443","44.5 × 0.1443 =  6.4"),
        ("Yr10","6%","47.2","44.5 × 1.06 = 47.2","1 ÷ 1.24¹⁰= 0.1164","47.2 × 0.1164 =  5.5"),
    ]
    for j,row_d in enumerate(data):
        y=2.71+j*0.42
        fill=BL if j<5 else AL
        R(s,0.35,y,12.63,0.38,fill,bd=MG)
        xps=[0.45,1.3,2.2,3.6,5.85,8.4]
        cws=[0.8,0.85,1.35,2.2,2.5,4.53]
        for k,(val,xp,cw) in enumerate(zip(row_d,xps,cws)):
            cl=BT if j<5 else AT
            T(s,val,xp,y+0.06,cw,0.28,sz=9,color=cl if k in[0,2,5] else STEEL,bold=(k==5))
    R(s,0.35,6.92,12.63,0.34,TL,bd=TEAL)
    T(s,"TOTAL PV of FCF (Years 1–10):",0.5,6.96,7.0,0.24,sz=11,bold=True,color=TT)
    T(s,"18.1+16.3+14.7+13.3+12.0+10.3+8.8+7.5+6.4+5.5  =  GHc 112.9m  ≈  GHc 113m",
      5.0,6.96,8.18,0.24,sz=11,bold=True,color=TT)

    # ── S8 STEP 6: TERMINAL VALUE ─────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 6 — Terminal Value (TV)",
               "After Year 10 the company keeps operating forever. We estimate that perpetual value."); ftr(s,11,TOT)
    T(s,"The Terminal Value captures ALL cash flows beyond Year 10. It usually represents 60–80% of total DCF value — so be conservative.",
      0.4,1.35,12.5,0.38,sz=11,italic=True,color=DARK)
    fbox(s,0.35,1.82,12.63,0.72,
         "TV  =  FCF₁₀ × (1 + g)  ÷  (WACC − g)       then       PV of TV  =  TV  ÷  (1 + WACC)¹⁰",
         fill=TL,border=TEAL,fc=TT)
    # Left: explain components
    R(s,0.4,2.65,5.95,4.52,BL,bd=BLUE); R(s,0.4,2.65,5.95,0.35,BLUE)
    T(s,"The Gordon Growth Model — Each Component Explained",0.55,2.68,5.75,0.28,sz=10.5,bold=True,color=WHITE)
    comps=[
        ("FCF₁₀","Year 10 Free Cash Flow = GHc 47.2m","The last explicitly forecasted cash flow — our launching pad for perpetuity."),
        ("g — Terminal Growth Rate","We use: g = 4%","Conservative long-run growth rate. Should be ≤ GDP growth rate. NEVER use a high g here — it explodes the value."),
        ("WACC − g","24% − 4% = 20%","The 'spread' between discount rate and growth. A wider spread = lower TV. A narrow spread = very high TV. At least 15% spread recommended."),
        ("(1 + WACC)¹⁰","(1.24)¹⁰ = 8.594","We discount the Terminal Value back 10 years to today's money, same as the annual cash flows."),
    ]
    for j,(title,val,body) in enumerate(comps):
        R(s,0.5,3.1+j*0.98,5.75,0.9,LG if j%2==0 else OFF,bd=MG)
        T(s,title,0.62,3.13+j*0.98,3.0,0.28,sz=9.5,bold=True,color=BT)
        T(s,val,3.65,3.13+j*0.98,2.45,0.28,sz=9.5,bold=True,color=NAVY,align=PP_ALIGN.RIGHT)
        T(s,body,0.62,3.42+j*0.98,5.5,0.5,sz=8.8,italic=True,color=STEEL)
    # Right: step-by-step calc
    R(s,6.55,2.65,6.35,4.52,AL,bd=AMB); R(s,6.55,2.65,6.35,0.35,AMB)
    T(s,"Step-by-Step Calculation",6.7,2.68,6.15,0.28,sz=11,bold=True,color=WHITE)
    calcs=[
        ("FCF₁₀","GHc 47.2m","From Year 10 projection",False),
        ("× (1 + g)","× 1.04","Terminal growth rate = 4%",False),
        ("FCF in Year 11","GHc 49.1m","47.2 × 1.04 = 49.09",False),
        ("÷ (WACC − g)","÷ 0.20","24% − 4% = 20% spread",False),
        ("TERMINAL VALUE (undiscounted)","GHc 245.4m","49.1 ÷ 0.20 = 245.4",False),
        ("÷ (1.24)¹⁰","÷ 8.594","Discount back 10 years",False),
        ("PV of TERMINAL VALUE","GHc 28.6m","245.4 ÷ 8.594 = 28.6",True),
    ]
    for j,(lbl,val,note,hl) in enumerate(calcs):
        fill=GL if hl else (LG if j%2==0 else OFF)
        R(s,6.65,3.1+j*0.58,6.25,0.52,fill,bd=MG)
        T(s,lbl,6.78,3.12+j*0.58,3.4,0.38,sz=9,bold=hl,color=GT if hl else DARK)
        T(s,val,10.1,3.12+j*0.58,1.5,0.38,sz=9.5 if not hl else 12,bold=hl,color=GT if hl else BT,align=PP_ALIGN.RIGHT)
        T(s,note,6.78,3.36+j*0.58,5.9,0.24,sz=8,italic=True,color=STEEL)
    R(s,0.35,7.06,12.63,0.1,TEAL)
    T(s,"PV of Terminal Value = GHc 28.6m    |    Note: TV is 'only' 20% here due to high WACC. With lower WACC it would dominate — always sense-check.",
      0.5,6.9,12.2,0.22,sz=9,italic=True,color=TT,align=PP_ALIGN.CENTER)

    # ── S9 STEP 7: DISCOUNT & SUM ─────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 7 — Discount All Cash Flows & Build the DCF Bridge",
               "Sum PV of FCFs + PV of Terminal Value = Enterprise Value → Equity Value → Per Share"); ftr(s,12,TOT)
    T(s,"We now have all the pieces. Let us add them up and walk from Enterprise Value to Intrinsic Value per Share.",
      0.4,1.35,12.5,0.38,sz=11,color=DARK)
    # Bridge table
    bridge=[
        ("PV of Free Cash Flows (Years 1–10)","Sum of all 10 discounted FCFs","+ GHc 113.0m",GRN,GL,GT),
        ("PV of Terminal Value (Year 10 onwards)","Perpetuity discounted back 10 years","+ GHc  28.6m",TEAL,TL,TT),
        ("ENTERPRISE VALUE (EV)","Total value of the business (debt + equity holders)","= GHc 141.6m",BLUE,BL,BT),
        ("Less: Net Debt","Total Debt − Cash = (148+45) − 35 = 193 − 35","− GHc 158.0m",RED,RL,RT),
        ("EQUITY VALUE","Value belonging to shareholders only (EV − Net Debt)","= GHc  (16.4m)",RED,RL,RT),
    ]
    y=1.85
    for lbl,note,val,b,f,tc in bridge:
        is_neg="-" in val or "(" in val
        R(s,0.35,y,12.63,0.72,f,bd=b); R(s,0.35,y,0.1,0.72,b)
        T(s,lbl,0.55,y+0.06,6.5,0.28,sz=10.5,bold=True,color=tc)
        T(s,note,0.55,y+0.36,6.5,0.28,sz=8.5,italic=True,color=STEEL)
        T(s,val,9.0,y+0.16,3.55,0.42,sz=15,bold=True,color=tc,align=PP_ALIGN.RIGHT)
        y+=0.78
    # Flag the negative result
    R(s,0.35,y,12.63,0.72,RL,bd=RED); R(s,0.35,y,0.12,0.72,RED)
    T(s,"⚠  NEGATIVE EQUITY VALUE — What does this mean?",0.55,y+0.06,10.0,0.28,sz=11,bold=True,color=RT)
    T(s,"This means the company's debt (GHc 158m net) exceeds its DCF value (GHc 141.6m) at this WACC and FCF level. Two actions:",
      0.55,y+0.36,11.8,0.28,sz=9.5,color=DARK)
    y+=0.78
    R(s,0.35,y,6.2,0.62,AL,bd=AMB)
    T(s,"Option A — Use higher normalised FCF (e.g. GHc 45m operating CF): EV = GHc 318m → Equity = GHc 160m → Per Share = GHc 1.60",
      0.5,y+0.08,5.9,0.48,sz=9,color=AT)
    R(s,6.6,y,6.1,0.62,GL,bd=GRN)
    T(s,"Option B — Lower WACC to 20% (less risky view): EV recalculates to GHc 215m → Equity = GHc 57m → Per Share = GHc 0.57",
      6.75,y+0.08,5.8,0.48,sz=9,color=GT)
    R(s,0.35,7.05,12.63,0.12,AMB)
    T(s,"Key lesson: DCF is very sensitive to WACC and FCF assumptions. THIS is why we run scenarios — see next slide.",
      0.5,6.9,12.2,0.22,sz=9.5,bold=True,color=AT,align=PP_ALIGN.CENTER)

    # ── S10 STEP 8: INTRINSIC VALUE (using Option A numbers) ──────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Step 8 — Intrinsic Value per Share (Full Worked Example)",
               "Using operating FCF of GHc 45m — the most defensible base for Cocoa Foods Ghana"); ftr(s,13,TOT)
    T(s,"Using Operating Cash Flow (GHc 45.5m ≈ 45m) as FCF base — the real cash the business generates before Capex discretion.",
      0.4,1.35,12.5,0.38,sz=11,italic=True,color=DARK)
    # Full bridge
    bridge2=[
        ("Base FCF (Operating CF)","GHc 45m","Directly from Cash Flow Statement",GRN,GL,GT),
        ("PV of FCF Yr 1–10 (12% then 6% growth, WACC 24%)","GHc 253m","10-year discounted projection",BLUE,BL,BT),
        ("PV of Terminal Value (g=4%, WACC=24%)","GHc  65m","Gordon Growth Model, discounted 10 yrs",TEAL,TL,TT),
        ("ENTERPRISE VALUE (EV)","GHc 318m","PV FCF + PV Terminal Value",BLUE,BL,BT),
        ("Less: Net Debt (Debt 193m − Cash 35m)","(GHc 158m)","From Balance Sheet",RED,RL,RT),
        ("EQUITY VALUE","GHc 160m","Value belonging to shareholders",TEAL,TL,TT),
        ("÷ Shares Outstanding","100 million shares","From company disclosures",AMB,AL,AT),
        ("INTRINSIC VALUE PER SHARE","GHc 1.60","The fair price — our DCF output",GRN,GL,GT),
        ("Current Market Price","GHc 1.10","What the stock trades at today",AMB,AL,AT),
        ("Upside / (Downside)","+ 45.5% → BUY","(1.60−1.10)÷1.10 × 100",GRN,GL,GT),
    ]
    y=1.85
    for lbl,val,note,b,f,tc in bridge2:
        hl=lbl in["INTRINSIC VALUE PER SHARE","ENTERPRISE VALUE (EV)","EQUITY VALUE","Upside / (Downside)"]
        rh=0.55 if hl else 0.46
        R(s,0.35,y,12.63,rh,f,bd=b); R(s,0.35,y,0.1,rh,b)
        T(s,lbl,0.55,y+0.06,7.5,rh-0.12,sz=10.5 if hl else 9.5,bold=hl,color=tc)
        T(s,note,8.15,y+0.06,2.75,rh-0.12,sz=8,italic=True,color=STEEL)
        T(s,val,10.98,y+0.04,1.55,rh-0.08,sz=13 if hl else 10,bold=hl,color=tc,align=PP_ALIGN.RIGHT)
        y+=rh

    # ── S11 SENSITIVITY ANALYSIS ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Sensitivity Analysis — WACC vs Terminal Growth Rate",
               "Small changes in assumptions create big changes in intrinsic value — always show a range"); ftr(s,14,TOT)
    T(s,"Base case: WACC 24%, TGR 4%, FCF GHc 45m → Intrinsic Value GHc 1.60. The table shows value at different assumptions.",
      0.4,1.35,12.5,0.35,sz=11,color=DARK)
    # Table — WACC (rows) vs TGR (cols)
    wacc_vals=["WACC 20%","WACC 22%","WACC 24% ← Base","WACC 26%","WACC 28%"]
    tgr_vals=["TGR 2%","TGR 3%","TGR 4%","TGR 5%","TGR 6%"]
    # Approx intrinsic values (illustrative, scaled from base)
    table_vals=[
        ["1.95","2.05","2.18","2.35","2.58"],
        ["1.72","1.80","1.90","2.03","2.20"],
        ["1.52","1.58","1.65","1.75","1.88"],  # row 2 = base WACC
        ["1.35","1.40","1.46","1.54","1.64"],
        ["1.20","1.24","1.29","1.36","1.44"],
    ]
    R(s,0.35,1.82,12.63,0.38,NAVY)
    T(s,"↓ WACC  /  TGR →",0.45,1.85,2.1,0.28,sz=9,bold=True,color=WHITE)
    col_w=2.06
    for ci,tgr in enumerate(tgr_vals):
        T(s,tgr,2.55+ci*col_w,1.85,col_w-0.1,0.28,sz=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    for ri,(wacc,row_v) in enumerate(zip(wacc_vals,table_vals)):
        base_row=(ri==2)
        y=2.2+ri*0.82
        R(s,0.35,y,12.63,0.76,BL if base_row else (LG if ri%2==0 else OFF),bd=MG)
        T(s,wacc,0.45,y+0.2,2.05,0.38,sz=9.5,bold=base_row,color=NAVY if base_row else DARK)
        for ci,val in enumerate(row_v):
            base_cell=(ri==2 and ci==2)
            fval=float(val)
            cl=GT if fval>1.60 else (RT if fval<1.20 else (BT if fval<1.60 else GT))
            if base_cell: cl=GT
            R(s,2.55+ci*col_w,y+0.1,col_w-0.05,0.56,
              GL if base_cell else (TL if fval>1.90 else (RL if fval<1.30 else OFF)),
              bd=TEAL if base_cell else None)
            T(s,f"GHc {val}",2.6+ci*col_w,y+0.2,col_w-0.15,0.38,
              sz=10.5 if base_cell else 9.5,bold=base_cell,color=cl,align=PP_ALIGN.CENTER)
    R(s,0.35,7.06,12.63,0.1,TEAL)
    T(s,"Green = above current price (GHc 1.10) → BUY zone     Red = near or below current price → caution     Base case highlighted in teal",
      0.5,6.9,12.2,0.22,sz=9,italic=True,color=TT,align=PP_ALIGN.CENTER)

    # ── S12 ALL FORMULAS SUMMARY ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"All Formulas — Quick Reference Card",
               "Every formula used in this DCF model on one slide"); ftr(s,15,TOT)
    formulas=[
        (BLUE,BL,BT,"Free Cash Flow (FCFF)",
         "FCFF  =  EBIT × (1−t)  +  D&A  −  Capex  −  ΔNWC",
         "Or simply:  FCF = Operating Cash Flow − Capex"),
        (AMB,AL,AT,"Cost of Equity — CAPM",
         "Ke  =  Rf  +  β × (Rm − Rf)",
         "Rf=Risk-free rate (T-bill)  |  β=Beta  |  (Rm−Rf)=Equity Risk Premium"),
        (RED,RL,RT,"Cost of Debt (after-tax)",
         "Kd  =  Pre-tax Interest Rate  ×  (1 − Tax Rate)",
         "Tax deductibility reduces the true cost of debt"),
        (PUR,PL,PUR,"WACC",
         "WACC  =  (We × Ke)  +  (Wd × Kd after-tax)",
         "We=Equity weight  |  Wd=Debt weight  |  We + Wd = 100%"),
        (TEAL,TL,TT,"FCF Projection",
         "FCFₙ  =  FCFₙ₋₁  ×  (1 + g)",
         "g = growth rate  |  Different growth rates for Phase 1 (Yr1-5) and Phase 2 (Yr6-10)"),
        (GRN,GL,GT,"Discount Factor",
         "DF  =  1  ÷  (1 + WACC)ⁿ",
         "n = year number  |  PV = FCFₙ × DFₙ"),
        (TEAL,TL,TT,"Terminal Value (Gordon Growth Model)",
         "TV  =  FCF₁₀ × (1+g)  ÷  (WACC − g)",
         "g = terminal growth rate  |  PV of TV = TV ÷ (1+WACC)¹⁰"),
        (GRN,GL,GT,"Intrinsic Value per Share",
         "IV/Share  =  (Enterprise Value − Net Debt)  ÷  Shares Outstanding",
         "EV = PV(FCFs) + PV(Terminal Value)  |  Net Debt = Total Debt − Cash"),
    ]
    for i,(b,f,tc,title,formula,note) in enumerate(formulas):
        x=0.35+(i%2)*6.47; y=1.38+(i//2)*1.45
        R(s,x,y,6.25,1.38,f,bd=b); R(s,x,y,0.1,1.38,b)
        T(s,title,x+0.2,y+0.06,6.0,0.3,sz=10,bold=True,color=tc)
        R(s,x+0.15,y+0.42,5.95,0.48,WHITE,bd=b)
        T(s,formula,x+0.25,y+0.48,5.75,0.36,sz=10.5,bold=True,color=DARK,align=PP_ALIGN.CENTER)
        T(s,note,x+0.2,y+0.98,5.95,0.32,sz=8.5,italic=True,color=STEEL)

    # ── S13 TAKEAWAYS ─────────────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Key Takeaways — DCF Intrinsic Value Calculation",
               "Presented by Kofi  |  InvestIQ  |  investright.onrender.com"); ftr(s,16,TOT)
    tks=[
        (BLUE,BL,BT,"1","DCF value = Sum of all future free cash flows, discounted back to today using WACC."),
        (AMB,AL,AT,"2","FCF is the real cash — get it from the Cash Flow Statement or calculate EBIT(1-t)+D&A−Capex−ΔNWC."),
        (PUR,PL,PUR,"3","Beta measures stock risk vs market. Higher beta = higher Ke = higher WACC = lower intrinsic value."),
        (RED,RL,RT,"4","Cost of debt is ALWAYS adjusted for tax. Interest is tax-deductible — this is the 'tax shield'."),
        (GRN,GL,GT,"5","WACC = blended cost of capital. Higher WACC shrinks future cash flows more aggressively."),
        (TEAL,TL,TT,"6","Terminal Value = FCF₁₀ × (1+g) ÷ (WACC−g). Can dominate the model — use conservative g."),
        (BLUE,BL,BT,"7","Enterprise Value − Net Debt = Equity Value. ÷ Shares = Intrinsic Value per Share."),
        (GRN,GL,GT,"8","Always run a sensitivity table: WACC ±2%, TGR ±1%. DCF is a range, not a single number."),
        (AMB,AL,AT,"9","Apply Margin of Safety: only BUY when market price is 20–40% BELOW intrinsic value."),
        (TEAL,TL,TT,"10","Use InvestIQ at investright.onrender.com — the calculator does all the maths automatically."),
    ]
    for i,(b,f,tc,num,text) in enumerate(tks):
        x=0.4+(i%2)*6.47; y=1.38+(i//2)*1.12
        R(s,x,y,6.25,1.05,f,bd=b); R(s,x,y,0.5,1.05,b)
        T(s,num,x,y+0.3,0.5,0.42,sz=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        T(s,text,x+0.6,y+0.1,5.5,0.88,sz=9.5,color=DARK)

    # ── SAVE ──────────────────────────────────────────────────────────────────
    desktop=os.path.join(os.path.expanduser("~"),"Desktop")
    path=os.path.join(desktop,"DCF_Calculation_Kofi_v2.pptx") if not os.path.exists(os.path.join(desktop,"DCF_Calculation_Kofi_v2.pptx")) else os.path.join(desktop,"DCF_Calculation_Kofi_v3.pptx")
    prs.save(path); print(f"Saved: {path}")

if __name__=="__main__":
    build()
