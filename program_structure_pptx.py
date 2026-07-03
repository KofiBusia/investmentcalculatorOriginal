"""
Financial Analysis & Valuation Program — Structure & Curriculum
Presenter: Kofi  |  InvestIQ
Saves: Desktop/Program_Structure_Kofi.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

WHITE=RGBColor(0xff,0xff,0xff); OFF=RGBColor(0xf8,0xfa,0xfc)
LG=RGBColor(0xf1,0xf5,0xf9); MG=RGBColor(0xe2,0xe8,0xf0)
STEEL=RGBColor(0x64,0x74,0x8b); DARK=RGBColor(0x0f,0x17,0x2a)
NAVY=RGBColor(0x1e,0x3a,0x8a); BLUE=RGBColor(0x1d,0x4e,0xd8)
BL=RGBColor(0xdb,0xea,0xff); BT=RGBColor(0x1e,0x40,0xaf)
TEAL=RGBColor(0x06,0x7a,0x5f); TL=RGBColor(0xd1,0xfa,0xee); TT=RGBColor(0x06,0x5f,0x46)
AMB=RGBColor(0xd9,0x77,0x06); AL=RGBColor(0xff,0xf7,0xd6); AT=RGBColor(0x92,0x40,0x02)
RED=RGBColor(0xb9,0x1c,0x1c); RL=RGBColor(0xfe,0xe2,0xe2); RT=RGBColor(0x99,0x1b,0x1b)
GRN=RGBColor(0x15,0x80,0x3d); GL=RGBColor(0xdc,0xfc,0xe8); GT=RGBColor(0x16,0x65,0x34)
PUR=RGBColor(0x6d,0x28,0xd9); PL=RGBColor(0xed,0xe9,0xfe); PT=RGBColor(0x4c,0x1d,0x95)
SUB=RGBColor(0xba,0xd0,0xf8)

def R(sl,l,t,w,h,fill,bd=None,bw=0.75):
    s=sl.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if bd: s.line.color.rgb=bd; s.line.width=Pt(bw)
    else: s.line.fill.background()

def T(sl,text,l,t,w,h,sz=10,bold=False,color=DARK,align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Calibri"

def hdr(sl,title,sub=None,ac=NAVY):
    R(sl,0,0,13.33,1.25,ac)
    T(sl,title,0.35,0.1,12.5,0.65,sz=24,bold=True,color=WHITE)
    if sub: T(sl,sub,0.35,0.82,12.5,0.38,sz=10,color=SUB,italic=True)

def ftr(sl,n,tot):
    R(sl,0,7.2,13.33,0.3,MG)
    T(sl,"Financial Analysis & Valuation Program  |  Presented by Kofi  |  InvestIQ  |  investright.onrender.com",
      0.25,7.22,11.5,0.24,sz=7.5,color=STEEL)
    T(sl,f"{n}/{tot}",12.5,7.22,0.7,0.24,sz=7.5,color=STEEL,align=PP_ALIGN.RIGHT)

def sbg(sl): R(sl,0,0,13.33,7.5,OFF)

def session_card(sl,x,y,w,h,num,title,topics,border,light,tc):
    R(sl,x,y,w,h,light,bd=border)
    R(sl,x,y,w,0.32,border)
    T(sl,f"Session {num}",x+0.12,y+0.06,w-0.24,0.2,sz=8.5,bold=True,color=WHITE)
    T(sl,title,x+0.12,y+0.4,w-0.24,0.32,sz=10,bold=True,color=tc)
    T(sl,topics,x+0.12,y+0.76,w-0.24,h-0.84,sz=8.5,color=DARK)

def build():
    prs=Presentation()
    prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    TOT=10

    # ── S1 COVER ──────────────────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    R(s,0,0,13.33,7.5,NAVY)
    R(s,0,0,13.33,0.22,AMB); R(s,0,7.28,13.33,0.22,TEAL); R(s,0,0.22,0.18,7.06,BLUE)
    T(s,"Financial Analysis",0.55,0.55,12.3,0.9,sz=42,bold=True,color=WHITE)
    T(s,"& Valuation Program",0.55,1.48,12.3,0.9,sz=42,bold=True,color=AMB)
    T(s,"From Reading Financial Statements to Calculating Intrinsic Value Per Share",
      0.55,2.52,12.3,0.45,sz=15,italic=True,color=SUB,align=PP_ALIGN.LEFT)
    R(s,0.55,3.18,12.25,2.65,RGBColor(0x1a,0x30,0x6e))
    R(s,0.55,3.18,12.25,0.06,AMB)
    cols=[
        (TEAL,TL,TT,"Stage 1","Financial Statement\nAnalysis","4 Weeks · 12 Sessions"),
        (AMB,AL,AT,"Stage 2","DCF Valuation\nFundamentals","Coming Soon"),
        (PUR,PL,PT,"Stage 3","Advanced Valuation\n& Portfolio","Coming Soon"),
    ]
    for i,(b,f,tc,stage,title,sub) in enumerate(cols):
        cx=0.75+i*4.1
        R(s,cx,3.32,3.9,2.35,RGBColor(0x14,0x26,0x5e),bd=b)
        R(s,cx,3.32,3.9,0.06,b)
        T(s,stage,cx+0.15,3.45,3.6,0.28,sz=9,bold=True,color=b)
        T(s,title,cx+0.15,3.78,3.6,0.62,sz=14,bold=True,color=WHITE)
        T(s,sub,cx+0.15,4.45,3.6,0.32,sz=9.5,color=RGBColor(0x94,0xa3,0xb8))
    T(s,"Presented by Kofi  ·  Powered by InvestIQ  ·  investright.onrender.com",
      0.4,7.0,12.5,0.22,sz=9,color=RGBColor(0x94,0xa3,0xb8),align=PP_ALIGN.CENTER)

    # ── S2 PROGRAM MISSION & OVERVIEW ─────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Program Mission & Learning Journey",
               "What you will be able to do by the end — and how we get there"); ftr(s,2,TOT)
    # Mission box
    R(s,0.35,1.32,12.63,0.88,BL,bd=BLUE)
    R(s,0.35,1.32,0.14,0.88,BLUE)
    T(s,"MISSION:",0.58,1.38,1.4,0.28,sz=11,bold=True,color=BT)
    T(s,"By the end of this program, every student will be able to open any company's annual report, understand every number in it, build a DCF valuation model from scratch, and arrive at a defensible intrinsic value per share — just like a professional investment analyst.",
      2.0,1.38,10.85,0.72,sz=10.5,color=DARK)
    # Three outcome columns
    cards=[
        (TEAL,TL,TT,"Stage 1 — Financial Analysis",
         ["Read and interpret an Income Statement, Balance Sheet, and Cash Flow Statement",
          "Understand how the 3 statements connect and feed into each other",
          "Calculate key financial ratios (profitability, liquidity, leverage, efficiency)",
          "Identify red flags and quality signals in company accounts",
          "Apply analysis to both a Manufacturing company and a Bank"],
         "4 Weeks  ·  12 Sessions  ·  NOW OPEN"),
        (AMB,AL,AT,"Stage 2 — DCF Valuation",
         ["Calculate Free Cash Flow (FCF, FCFF, FCFE) from financial statements",
          "Compute Beta using regression and the Hamada equation",
          "Build Cost of Equity (CAPM) and Cost of Debt calculations",
          "Construct a full WACC from a company's capital structure",
          "Project FCF, compute Terminal Value, discount to intrinsic value"],
         "Follows Stage 1  ·  Opening Soon"),
        (PUR,PL,PT,"Stage 3 — Advanced Valuation",
         ["Comparable company analysis (Trading Comps)",
          "Precedent transaction analysis (Deal Comps)",
          "Relative valuation: P/E, P/B, EV/EBITDA multiples",
          "Building a full investment thesis and buy/sell case",
          "Live stock pitch presentation"],
         "Advanced  ·  Date TBD"),
    ]
    for i,(b,f,tc,title,pts,tag) in enumerate(cards):
        cx=0.35+i*4.35
        R(s,cx,2.32,4.18,4.52,f,bd=b); R(s,cx,2.32,4.18,0.35,b)
        T(s,title,cx+0.15,2.35,4.0,0.26,sz=10.5,bold=True,color=WHITE)
        for j,pt in enumerate(pts):
            T(s,f"✓  {pt}",cx+0.15,2.78+j*0.58,3.9,0.48,sz=8.8,color=DARK)
        R(s,cx,6.6,4.18,0.2,b)
        T(s,tag,cx+0.15,6.62,3.9,0.16,sz=7.5,bold=True,color=WHITE)

    # ── S3 FULL PROGRAM SNAPSHOT ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Program at a Glance — All Three Stages",
               "The complete roadmap from financial literacy to professional valuation"); ftr(s,3,TOT)
    stages=[
        (TEAL,TL,TT,"STAGE 1","Financial Statement Analysis","4 Weeks  |  12 Sessions",
         [("Week 1","Foundations","Why financial statements matter · Introduction to the 3 statements · How they are structured"),
          ("Week 2","The Income Statement","Line-by-line breakdown · Manufacturing vs Bank · Key IS ratios"),
          ("Week 3","Balance Sheet & Cash Flow","BS line by line · CFS Operating/Investing/Financing · How all 3 connect"),
          ("Week 4","Ratio Analysis & Application","Profitability · Liquidity · Leverage · Efficiency · Reading an annual report"),
         ]),
        (AMB,AL,AT,"STAGE 2","DCF Valuation Fundamentals","4 Weeks  |  12 Sessions",
         [("Week 5","Free Cash Flow","FCF · FCFF · FCFE · Where to find components · Normalising FCF"),
          ("Week 6","Cost of Capital","Beta regression · Hamada equation · CAPM (Ke) · Cost of Debt (Kd)"),
          ("Week 7","WACC & Projections","Building WACC · Projecting FCF for 10 years · Growth rate assumptions"),
          ("Week 8","Intrinsic Value","Terminal Value · Discounting · EV → Equity Value → Price per Share"),
         ]),
        (PUR,PL,PT,"STAGE 3","Advanced Valuation & Portfolio","4 Weeks  |  12 Sessions",
         [("Week 9","Relative Valuation","P/E · P/B · EV/EBITDA · How to find peer multiples · When to use comps"),
          ("Week 10","Transaction Analysis","Precedent deals · Control premium · Synergies · Deal structure basics"),
          ("Week 11","Investment Thesis","Bull/Bear case · Catalysts · Risk factors · Margin of safety"),
          ("Week 12","Live Stock Pitch","Present a full investment recommendation using InvestIQ calculator"),
         ]),
    ]
    for i,(b,f,tc,stage_lbl,stage_title,meta,weeks) in enumerate(stages):
        sy=1.32+i*1.98
        R(s,0.35,sy,12.63,1.88,f,bd=b)
        R(s,0.35,sy,1.45,1.88,b)
        T(s,stage_lbl,0.35,sy+0.52,1.45,0.4,sz=9.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
        T(s,stage_title,1.9,sy+0.08,4.5,0.4,sz=13,bold=True,color=tc)
        T(s,meta,1.9,sy+0.52,4.5,0.28,sz=9,color=STEEL,italic=True)
        for j,(wk,wt,desc) in enumerate(weeks):
            wx=6.55+j*1.68
            R(s,wx,sy+0.12,1.58,1.62,WHITE,bd=b)
            T(s,wk,wx+0.08,sy+0.16,1.42,0.22,sz=7.5,bold=True,color=b)
            T(s,wt,wx+0.08,sy+0.4,1.42,0.26,sz=9,bold=True,color=tc)
            T(s,desc,wx+0.08,sy+0.7,1.42,0.98,sz=7.5,color=DARK)

    # ── S4 STAGE 1 AT A GLANCE ────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Financial Statement Analysis at a Glance",
               "4 weeks · 3 sessions per week · 12 total sessions · InvestIQ platform throughout"); ftr(s,4,TOT)
    # Stats bar
    stats=[("4","Weeks"),("3","Sessions / Week"),("12","Total Sessions"),("2","Case Companies"),("100+","Financial Ratios")]
    for i,(num,lbl) in enumerate(stats):
        cx=0.35+i*2.58
        R(s,cx,1.32,2.44,0.82,BL,bd=BLUE)
        T(s,num,cx+0.12,1.34,1.1,0.48,sz=28,bold=True,color=BLUE)
        T(s,lbl,cx+1.22,1.46,1.1,0.3,sz=9,color=BT)
    # Four week overview
    weeks_data=[
        (TEAL,TL,TT,"Week 1","Foundations & Why It Matters",
         "Sessions 1–3",
         "• Why do financial statements exist?\n• Who reads them and why it matters to investors\n• Introduction to the 3 core statements\n• The accounting equation: Assets = Liabilities + Equity\n• Overview of both case companies"),
        (BLUE,BL,BT,"Week 2","The Income Statement in Depth",
         "Sessions 4–6",
         "• IS line-by-line: Revenue → COGS → Gross Profit → EBIT → Net Profit\n• Manufacturing IS (Cocoa Foods Ghana) — full worked example\n• Bank IS (GoldCoast Bank) — how it differs\n• Key IS ratios: Gross Margin, Net Margin, EBIT Margin\n• Common accounting adjustments"),
        (AMB,AL,AT,"Week 3","Balance Sheet & Cash Flow Statement",
         "Sessions 7–9",
         "• Balance Sheet: Assets, Liabilities, Equity — line by line\n• BS for Manufacturing vs Bank — key structural differences\n• The Cash Flow Statement: Operating, Investing, Financing\n• How to trace Net Profit → CFS → Balance Sheet\n• Working capital and why it matters"),
        (PUR,PL,PT,"Week 4","Ratios, Connections & Application",
         "Sessions 10–12",
         "• Profitability ratios: ROE, ROA, ROCE\n• Liquidity ratios: Current Ratio, Quick Ratio\n• Leverage: Debt/Equity, Interest Coverage, Net Debt\n• Efficiency: Asset Turnover, Receivable Days, Payable Days\n• How all 3 statements feed into each other\n• Reading a real annual report — red flags checklist\n• Stage 1 Assessment"),
    ]
    for i,(b,f,tc,wk,title,meta,body) in enumerate(weeks_data):
        cx=0.35+(i%2)*6.47; cy=2.26+(i//2)*2.42
        R(s,cx,cy,6.25,2.32,f,bd=b); R(s,cx,cy,6.25,0.35,b)
        T(s,f"{wk}  —  {title}",cx+0.15,cy+0.06,5.9,0.22,sz=11,bold=True,color=WHITE)
        T(s,meta,cx+0.15,cy+0.3,5.9,0.18,sz=8.5,italic=True,color=WHITE)
        T(s,body,cx+0.15,cy+0.54,5.9,1.62,sz=9,color=DARK)

    # ── S5 WEEK 1 — SESSIONS 1–3 ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Week 1: Foundations",
               "Why Financial Statements Exist and What They Tell Us  |  Sessions 1 · 2 · 3"); ftr(s,5,TOT)
    R(s,0.35,1.32,12.63,0.38,TL,bd=TEAL)
    T(s,"WEEK 1 GOAL:  Students understand WHY financial statements exist, who uses them, and can identify the 3 core statements and their purpose.",
      0.5,1.38,12.2,0.26,sz=10,bold=True,color=TT)
    sw=4.12; gap=0.1
    sessions_w1=[
        (TEAL,TL,TT,"1",
         "Why Financial Statements Matter",
         "• The language of business — what a company is legally required to publish\n• Who reads financial statements: investors, banks, regulators, management\n• The 3 financial statements: IS, BS, CFS — what each answers\n• Introducing our two case companies:\n   – Cocoa Foods Ghana Ltd (Manufacturing)\n   – GoldCoast Bank Ltd (Banking)\n• Walk through the company profiles: what they do, where they operate, their scale\n• How an investor uses these statements to decide: Buy, Hold, or Sell?",
         "Whiteboard overview\nCompany profile handout\nInvestIQ platform tour"),
        (BLUE,BL,BT,"2",
         "The Accounting Equation & Structure of Statements",
         "• The accounting equation: Assets = Liabilities + Equity\n• Double-entry principle (simplified): every transaction affects 2 things\n• Structure of the Income Statement: top-down (Revenue → Net Profit)\n• Structure of the Balance Sheet: snapshot on a specific date\n• Structure of the Cash Flow Statement: cash movements in a period\n• Accrual accounting vs cash accounting — why they differ\n• Why Net Profit ≠ Cash in the bank",
         "Accounting equation worksheet\nSide-by-side statement template\nReal-world analogy: household budget"),
        (AMB,AL,AT,"3",
         "Reading Your First Financial Statement",
         "• Guided walk through Cocoa Foods Ghana's full IS\n• Identify each line: Revenue, COGS, Gross Profit, D&A, EBIT, Interest, Tax, Net Profit\n• What does each line tell us about the business?\n• Class exercise: fill in missing line items from narrative\n• Introduction to GoldCoast Bank's IS — why it looks so different\n• First ratio: Gross Margin = Gross Profit ÷ Revenue\n• Week 1 recap and preview of Week 2",
         "Printed IS handout (Cocoa Foods)\nFirst ratio exercise\nClass Q&A"),
    ]
    for i,(b,f,tc,num,title,body,resources) in enumerate(sessions_w1):
        cx=0.35+i*(sw+gap)
        R(s,cx,1.82,sw,5.44,f,bd=b); R(s,cx,1.82,sw,0.35,b)
        T(s,f"Session {num}",cx+0.14,1.85,sw-0.28,0.22,sz=9.5,bold=True,color=WHITE)
        T(s,title,cx+0.14,2.24,sw-0.28,0.36,sz=11,bold=True,color=tc)
        T(s,body,cx+0.14,2.66,sw-0.28,3.5,sz=9,color=DARK)
        R(s,cx+0.1,6.24,sw-0.2,0.94,WHITE,bd=b)
        T(s,"Resources / Activities",cx+0.2,6.28,sw-0.4,0.2,sz=7.5,bold=True,color=tc)
        T(s,resources,cx+0.2,6.5,sw-0.4,0.64,sz=8,italic=True,color=STEEL)

    # ── S6 WEEK 2 — SESSIONS 4–6 ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Week 2: The Income Statement",
               "Line-by-Line Breakdown for Manufacturing and Banking Companies  |  Sessions 4 · 5 · 6"); ftr(s,6,TOT)
    R(s,0.35,1.32,12.63,0.38,BL,bd=BLUE)
    T(s,"WEEK 2 GOAL:  Students can read any Income Statement line by line, explain each item, and calculate profitability ratios from scratch.",
      0.5,1.38,12.2,0.26,sz=10,bold=True,color=BT)
    sessions_w2=[
        (BLUE,BL,BT,"4",
         "Manufacturing IS — Deep Dive",
         "• Full IS walkthrough: Cocoa Foods Ghana Ltd\n• Revenue recognition — when is revenue recorded?\n• Cost of Goods Sold (COGS) — what goes in here?\n• Gross Profit and Gross Margin — first profitability signal\n• Operating expenses: Sales, Admin, Distribution\n• Depreciation & Amortisation — non-cash but important\n• EBIT (Earnings Before Interest and Tax) — why analysts love it\n• Finance costs: Interest expense and where it comes from",
         "Full IS printout\nGross margin calculation drill\nExcel template starter"),
        (TEAL,TL,TT,"5",
         "Manufacturing IS Continued + Bank IS",
         "• Tax calculation: Effective tax rate vs statutory rate\n• Net Profit After Tax — the 'bottom line'\n• IS Ratios: Gross Margin, EBIT Margin, Net Margin, Tax Rate\n• Class exercise: calculate all 4 margins for Cocoa Foods\n• Introduction to GoldCoast Bank IS:\n   – Interest Income vs Interest Expense\n   – Net Interest Income (NII)\n   – Loan Loss Provisions — what they are and why they spike in crises\n   – Non-interest income: fees, commissions, forex\n• How a Bank IS differs structurally from a Manufacturing IS",
         "Ratio calculation worksheet\nBank IS handout\nComparison table: Manufacturing vs Bank"),
        (AMB,AL,AT,"6",
         "IS Analysis & Quality of Earnings",
         "• Quality of earnings: sustainable vs one-off income\n• Red flags in the IS: declining margins, revenue without cash, rising provisions\n• Green flags: expanding margins, consistent revenue growth, growing EBIT\n• How IS connects to the Balance Sheet: retained earnings\n• How IS connects to the CFS: Net Profit is the starting point\n• Class exercise: analyse a 3-year IS trend for Cocoa Foods\n• Introduction to earnings manipulation — what to watch for\n• Week 2 quiz and recap",
         "3-year trend analysis sheet\nRed flag checklist\nMini quiz (5 questions)"),
    ]
    for i,(b,f,tc,num,title,body,resources) in enumerate(sessions_w2):
        cx=0.35+i*(sw+gap)
        R(s,cx,1.82,sw,5.44,f,bd=b); R(s,cx,1.82,sw,0.35,b)
        T(s,f"Session {num}",cx+0.14,1.85,sw-0.28,0.22,sz=9.5,bold=True,color=WHITE)
        T(s,title,cx+0.14,2.24,sw-0.28,0.36,sz=11,bold=True,color=tc)
        T(s,body,cx+0.14,2.66,sw-0.28,3.5,sz=9,color=DARK)
        R(s,cx+0.1,6.24,sw-0.2,0.94,WHITE,bd=b)
        T(s,"Resources / Activities",cx+0.2,6.28,sw-0.4,0.2,sz=7.5,bold=True,color=tc)
        T(s,resources,cx+0.2,6.5,sw-0.4,0.64,sz=8,italic=True,color=STEEL)

    # ── S7 WEEK 3 — SESSIONS 7–9 ──────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Week 3: Balance Sheet & Cash Flow Statement",
               "Assets, Liabilities, Equity and the Movement of Cash  |  Sessions 7 · 8 · 9"); ftr(s,7,TOT)
    R(s,0.35,1.32,12.63,0.38,GL,bd=GRN)
    T(s,"WEEK 3 GOAL:  Students can read both the Balance Sheet and Cash Flow Statement, and trace every major number between all three statements.",
      0.5,1.38,12.2,0.26,sz=10,bold=True,color=GT)
    sessions_w3=[
        (GRN,GL,GT,"7",
         "The Balance Sheet — Line by Line",
         "• What the BS shows: financial position on ONE date\n• Assets: Current (Cash, Receivables, Inventory) vs Non-Current (PP&E, Intangibles)\n• Liabilities: Current (Payables, Short-term Debt) vs Long-term (Bonds, LT Debt)\n• Equity: Share Capital + Retained Earnings + Reserves\n• Cocoa Foods Ghana BS — walked through in full\n• GoldCoast Bank BS — Loans as main asset, Deposits as main liability\n• The accounting equation in action: checking it balances\n• Key BS metric: Net Debt = Total Debt − Cash",
         "Full BS printout (both companies)\nAccounting equation checker\nNet debt calculation"),
        (TEAL,TL,TT,"8",
         "The Cash Flow Statement — Line by Line",
         "• Why CFS exists: profit ≠ cash\n• Section 1 — Operating Activities:\n   Start with Net Profit, add back D&A, adjust working capital\n• Section 2 — Investing Activities:\n   Capex, acquisitions, proceeds from asset sales\n• Section 3 — Financing Activities:\n   New debt raised, debt repaid, dividends paid, shares issued\n• Cocoa Foods Ghana CFS — walked through in full\n• GoldCoast Bank CFS — unique items (loan book growth as CFI)\n• Closing cash balance must tie to BS Cash line\n• Free Cash Flow = Operating CF − Capex (first sight of FCF)",
         "Full CFS printout (both companies)\nCash reconciliation exercise\nFCF calculation worksheet"),
        (BLUE,BL,BT,"9",
         "How the 3 Statements Connect",
         "• The 3-statement model: IS → BS → CFS\n• Trace: Net Profit (IS) → Retained Earnings (BS equity)\n• Trace: Capex (CFS Investing) → PP&E increase (BS assets)\n• Trace: Depreciation (IS) → PP&E decrease (BS) → CFS add-back\n• Trace: New debt raised (CFS Financing) → Debt on BS\n• Trace: Interest expense (IS) → paid in CFS Operating\n• Working Capital: (Current Assets − Cash) − Current Liabilities\n• Class exercise: complete a 3-statement linkage map\n• Common exam/interview questions on statement linkages",
         "3-statement linkage map template\nFlow diagram handout\nInterview prep Q&A sheet"),
    ]
    for i,(b,f,tc,num,title,body,resources) in enumerate(sessions_w3):
        cx=0.35+i*(sw+gap)
        R(s,cx,1.82,sw,5.44,f,bd=b); R(s,cx,1.82,sw,0.35,b)
        T(s,f"Session {num}",cx+0.14,1.85,sw-0.28,0.22,sz=9.5,bold=True,color=WHITE)
        T(s,title,cx+0.14,2.24,sw-0.28,0.36,sz=11,bold=True,color=tc)
        T(s,body,cx+0.14,2.66,sw-0.28,3.5,sz=9,color=DARK)
        R(s,cx+0.1,6.24,sw-0.2,0.94,WHITE,bd=b)
        T(s,"Resources / Activities",cx+0.2,6.28,sw-0.4,0.2,sz=7.5,bold=True,color=tc)
        T(s,resources,cx+0.2,6.5,sw-0.4,0.64,sz=8,italic=True,color=STEEL)

    # ── S8 WEEK 4 — SESSIONS 10–12 ────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Week 4: Financial Ratios & Applied Analysis",
               "Profitability · Liquidity · Leverage · Efficiency · Reading Annual Reports  |  Sessions 10 · 11 · 12"); ftr(s,8,TOT)
    R(s,0.35,1.32,12.63,0.38,AL,bd=AMB)
    T(s,"WEEK 4 GOAL:  Students can calculate all major financial ratios, interpret them in context, read a real annual report and spot red flags.",
      0.5,1.38,12.2,0.26,sz=10,bold=True,color=AT)
    sessions_w4=[
        (AMB,AL,AT,"10",
         "Financial Ratios — The Full Toolkit",
         "PROFITABILITY:\n• Gross Margin, EBIT Margin, Net Margin\n• Return on Equity (ROE) = Net Profit ÷ Equity\n• Return on Assets (ROA) = Net Profit ÷ Total Assets\n• Return on Capital Employed (ROCE)\n\nLIQUIDITY:\n• Current Ratio = Current Assets ÷ Current Liabilities\n• Quick Ratio = (Current Assets − Inventory) ÷ CL\n• Cash Ratio\n\nLEVERAGE:\n• Debt/Equity, Net Debt/EBITDA, Interest Coverage",
         "Full ratio formula sheet\nCocoa Foods ratio calculator\nBenchmark comparison table"),
        (RED,RL,RT,"11",
         "Efficiency Ratios & Reading Annual Reports",
         "EFFICIENCY RATIOS:\n• Asset Turnover = Revenue ÷ Total Assets\n• Receivable Days = Receivables ÷ Revenue × 365\n• Payable Days = Payables ÷ COGS × 365\n• Inventory Days = Inventory ÷ COGS × 365\n• Cash Conversion Cycle = Receivable + Inventory − Payable Days\n\nREADING ANNUAL REPORTS:\n• Where to find each statement in an annual report\n• Notes to the accounts — why they matter\n• Auditor's report — what 'qualified' means\n• Red flags: related-party transactions, aggressive revenue, growing receivables\n• Green flags: improving cash conversion, rising ROCE, declining debt",
         "Annual report extract handout\nRed flag checklist\nEfficiency ratio worksheet"),
        (GRN,GL,GT,"12",
         "Stage 1 Assessment & Capstone",
         "SESSION FORMAT (90 min):\n\n• Part 1 (30 min) — Written Assessment:\n   Read 2 pages of financial statements. Answer 10 questions covering IS ratios, BS structure, CFS analysis and statement linkages.\n\n• Part 2 (30 min) — Ratio Analysis:\n   Given Cocoa Foods Ghana full statements, calculate 8 key ratios and interpret each one.\n\n• Part 3 (30 min) — Group Presentations:\n   Each group presents one key finding from the case company and one red flag or green flag they identified.\n\n• Feedback & Stage 2 Preview:\n   What comes next — DCF Valuation begins next week.",
         "Assessment paper (printed)\nRatio answer sheet\nPresentation rubric\nStage 2 preview handout"),
    ]
    for i,(b,f,tc,num,title,body,resources) in enumerate(sessions_w4):
        cx=0.35+i*(sw+gap)
        R(s,cx,1.82,sw,5.44,f,bd=b); R(s,cx,1.82,sw,0.35,b)
        T(s,f"Session {num}",cx+0.14,1.85,sw-0.28,0.22,sz=9.5,bold=True,color=WHITE)
        T(s,title,cx+0.14,2.24,sw-0.28,0.36,sz=11,bold=True,color=tc)
        T(s,body,cx+0.14,2.66,sw-0.28,3.5,sz=9,color=DARK)
        R(s,cx+0.1,6.24,sw-0.2,0.94,WHITE,bd=b)
        T(s,"Resources / Activities",cx+0.2,6.28,sw-0.4,0.2,sz=7.5,bold=True,color=tc)
        T(s,resources,cx+0.2,6.5,sw-0.4,0.64,sz=8,italic=True,color=STEEL)

    # ── S9 LEARNING OUTCOMES & COMPETENCY MAP ────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Stage 1 — Learning Outcomes & Competency Map",
               "What students can DO by the end of each week — measurable, observable skills"); ftr(s,9,TOT)
    outcomes=[
        (TEAL,TL,TT,"After Week 1",
         [("Describe","the purpose of each of the 3 financial statements in one sentence"),
          ("Identify","the key sections and line items in an Income Statement"),
          ("Explain","why Net Profit does not equal cash in the bank"),
          ("Compare","the structure of a Manufacturing company vs a Bank at a high level"),
         ]),
        (BLUE,BL,BT,"After Week 2",
         [("Calculate","Gross Margin, EBIT Margin, Net Margin and Net Profit from raw IS data"),
          ("Explain","the role of Depreciation, Interest, and Tax in moving from Revenue to Net Profit"),
          ("Distinguish","between a Manufacturing IS and a Bank IS — key structural differences"),
          ("Spot","at least 3 red flags that suggest earnings quality may be poor"),
         ]),
        (GRN,GL,GT,"After Week 3",
         [("Read","a Balance Sheet and identify all major asset, liability and equity line items"),
          ("Trace","Net Profit from IS through to Retained Earnings on the Balance Sheet"),
          ("Explain","what each section of the Cash Flow Statement represents"),
          ("Produce","a 3-statement linkage diagram showing how IS, BS, and CFS connect"),
         ]),
        (AMB,AL,AT,"After Week 4",
         [("Calculate","ROE, ROA, Current Ratio, Quick Ratio, Debt/Equity, Asset Turnover, Receivable Days"),
          ("Interpret","each ratio in the context of the industry and the company's trend"),
          ("Navigate","a real annual report and locate any financial statement or note"),
          ("Present","a structured financial analysis of a company with a clear view on financial health"),
         ]),
    ]
    for i,(b,f,tc,label,items) in enumerate(outcomes):
        cx=0.35+(i%2)*6.47; cy=1.32+(i//2)*2.92
        R(s,cx,cy,6.25,2.82,f,bd=b); R(s,cx,cy,6.25,0.35,b)
        T(s,label,cx+0.15,cy+0.06,6.0,0.22,sz=11,bold=True,color=WHITE)
        T(s,"Students will be able to:",cx+0.15,cy+0.44,6.0,0.22,sz=9,italic=True,color=tc)
        for j,(verb,rest) in enumerate(items):
            ry=cy+0.7+j*0.5
            T(s,verb,cx+0.15,ry,0.85,0.38,sz=9.5,bold=True,color=tc)
            T(s,rest,cx+1.02,ry,5.0,0.38,sz=9.5,color=DARK)

    # ── S10 TOOLS & NEXT STEPS ────────────────────────────────────────────────
    s=prs.slides.add_slide(blank)
    sbg(s); hdr(s,"Tools, Resources & What Comes Next",
               "Everything students need to succeed in Stage 1 — and the path ahead"); ftr(s,10,TOT)
    # Tools
    R(s,0.35,1.32,7.9,3.52,BL,bd=BLUE); R(s,0.35,1.32,7.9,0.35,BLUE)
    T(s,"Tools & Resources Used in This Program",0.5,1.35,7.7,0.26,sz=11,bold=True,color=WHITE)
    tools=[
        ("InvestIQ Platform","investright.onrender.com","Live DCF calculator, financial analysis tools. Used from Stage 1 for ratio analysis and from Stage 2 for full valuation."),
        ("Presentation Decks","Slides provided each session","Financial Statements deck (Manufacturing & Bank) and DCF Calculation deck — pre-built, fully worked."),
        ("Case Study Companies","Cocoa Foods Ghana + GoldCoast Bank","Real-structure illustrative financials. All numbers realistic — balance sheets balance, cash flows tie."),
        ("Excel Templates","Provided Week 1","Ratio calculator, 3-statement model template, FCF workbook. Students fill them in during sessions."),
    ]
    for j,(tool,loc,desc) in enumerate(tools):
        R(s,0.42,1.76+j*0.76,7.75,0.7,LG if j%2==0 else OFF,bd=MG)
        T(s,tool,0.55,1.8+j*0.76,2.2,0.28,sz=9.5,bold=True,color=BT)
        T(s,loc,2.8,1.8+j*0.76,1.7,0.28,sz=8,italic=True,color=STEEL)
        T(s,desc,4.55,1.8+j*0.76,3.45,0.56,sz=8.8,color=DARK)
    # What comes next
    R(s,8.4,1.32,4.78,3.52,GL,bd=GRN); R(s,8.4,1.32,4.78,0.35,GRN)
    T(s,"After Stage 1 — What's Next",8.55,1.35,4.58,0.26,sz=11,bold=True,color=WHITE)
    nxt=[("Stage 2 opens","immediately after Stage 1 Assessment"),
         ("Topic: Free Cash Flow","FCF, FCFF, FCFE from statements"),
         ("Topic: Beta","Regression, Hamada, CAPM"),
         ("Topic: WACC","Cost of equity + cost of debt"),
         ("Topic: DCF Model","Full 10-year projection + Terminal Value"),
         ("Topic: Intrinsic Value","EV → Equity Value → Price per Share"),
         ("Platform","Full InvestIQ DCF calculator unlocked")]
    for j,(k,v) in enumerate(nxt):
        R(s,8.47,1.76+j*0.44,4.64,0.38,GL if j%2==0 else OFF,bd=MG)
        T(s,k,8.57,1.8+j*0.44,1.65,0.26,sz=8.5,bold=True,color=GT)
        T(s,v,10.28,1.8+j*0.44,2.68,0.26,sz=8.5,color=DARK)
    # Enrollment / contact
    R(s,0.35,4.96,12.63,2.18,NAVY)
    T(s,"Enroll Now — Stage 1 is Open",0.5,5.04,12.2,0.42,sz=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,"investright.onrender.com",0.5,5.52,12.2,0.38,sz=14,color=AMB,align=PP_ALIGN.CENTER)
    T(s,"Questions? Contact Kofi directly  ·  Program delivered live (in-person or virtual)  ·  All slides and materials provided  ·  Certificate on completion of Stage 1 assessment",
      0.5,5.96,12.2,0.28,sz=9.5,italic=True,color=SUB,align=PP_ALIGN.CENTER)
    T(s,"Powered by InvestIQ  ·  Crestline Technologies",
      0.5,6.44,12.2,0.28,sz=9,color=RGBColor(0x94,0xa3,0xb8),align=PP_ALIGN.CENTER)

    # ── SAVE ──────────────────────────────────────────────────────────────────
    desktop=os.path.join(os.path.expanduser("~"),"Desktop")
    path=os.path.join(desktop,"Program_Structure_Kofi.pptx")
    prs.save(path); print(f"Saved: {path}")

if __name__=="__main__":
    build()
