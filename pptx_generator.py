"""
YIN Stock Pitch Competition — 30-Slide PowerPoint Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0f, 0x17, 0x2a)   # dark navy bg
BLUE      = RGBColor(0x1d, 0x4e, 0xd8)   # primary blue
LIGHTBLUE = RGBColor(0x60, 0xa5, 0xfa)   # accent
GOLD      = RGBColor(0xf5, 0x9e, 0x0b)   # YIN gold/amber
RED       = RGBColor(0xdc, 0x26, 0x26)   # YIN red accent
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY      = RGBColor(0x64, 0x74, 0x8b)
LIGHTGRAY = RGBColor(0xf1, 0xf5, 0xf9)
GREEN     = RGBColor(0x05, 0x96, 0x69)
DARKGRAY  = RGBColor(0x1e, 0x29, 0x3b)


def _rgb(color):
    return color


def _add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, l, t, w, h, font_size=14, bold=False,
              color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def _dark_slide(slide):
    """Full dark navy background."""
    _add_rect(slide, 0, 0, 13.33, 7.5, NAVY)


def _header_bar(slide, title, subtitle=None, accent=BLUE):
    """Top accent bar + title."""
    _add_rect(slide, 0, 0, 13.33, 1.35, accent)
    _add_text(slide, title, 0.35, 0.08, 12, 0.7,
              font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_text(slide, subtitle, 0.35, 0.82, 12, 0.45,
                  font_size=11, color=RGBColor(0xba, 0xd0, 0xf8), align=PP_ALIGN.LEFT)


def _footer(slide, slide_num, total=30, label="YIN Stock Pitch Competition"):
    _add_rect(slide, 0, 7.15, 13.33, 0.35, DARKGRAY)
    _add_text(slide, label, 0.2, 7.17, 9, 0.28, font_size=8,
              color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.LEFT)
    _add_text(slide, f"{slide_num} / {total}", 12.5, 7.17, 0.7, 0.28,
              font_size=8, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.RIGHT)


def _bullet_box(slide, items, l, t, w, h, title=None,
                bg=DARKGRAY, title_color=LIGHTBLUE, text_color=WHITE,
                font_size=11, title_size=12):
    _add_rect(slide, l, t, w, h, bg)
    y = t + 0.1
    if title:
        _add_text(slide, title, l + 0.15, y, w - 0.3, 0.32,
                  font_size=title_size, bold=True, color=title_color)
        y += 0.35
    for item in items:
        _add_text(slide, f"  •  {item}", l + 0.1, y, w - 0.25, 0.3,
                  font_size=font_size, color=text_color)
        y += 0.28
    return y


def _placeholder_box(slide, label, l, t, w, h,
                     bg=RGBColor(0x1e, 0x29, 0x3b),
                     border=RGBColor(0x33, 0x41, 0x55)):
    """Dashed-border placeholder area."""
    _add_rect(slide, l, t, w, h, bg, border)
    _add_text(slide, f"[ {label} ]", l + 0.1, t + h / 2 - 0.2, w - 0.2, 0.4,
              font_size=10, italic=True,
              color=RGBColor(0x64, 0x74, 0x8b), align=PP_ALIGN.CENTER)


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # completely blank

    slides_data = []

    # ── SLIDE 1: Cover ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    # Gold accent strip at top
    _add_rect(s, 0, 0, 13.33, 0.18, GOLD)
    # Red accent bottom
    _add_rect(s, 0, 7.32, 13.33, 0.18, RED)
    # YIN badge
    _add_rect(s, 0.5, 0.6, 3.2, 1.1, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "YOUNG INVESTORS NETWORK", 0.55, 0.65, 3.1, 0.45,
              font_size=9, bold=True, color=GOLD)
    _add_text(s, "× InvestIQ", 0.55, 1.05, 3.1, 0.45,
              font_size=11, bold=True, color=LIGHTBLUE)
    # Main title
    _add_text(s, "STOCK PITCH", 0.5, 2.0, 12, 1.1,
              font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "COMPETITION", 0.5, 3.0, 12, 1.1,
              font_size=54, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _add_text(s, "A Complete Pitch Deck Template — Ghana Stock Exchange", 0.5, 4.15, 12, 0.5,
              font_size=16, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
    # Team info placeholders
    _add_rect(s, 2.5, 5.0, 8.33, 1.6, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Team Name:", 2.7, 5.1, 3.5, 0.35, font_size=11, color=GRAY)
    _add_text(s, "[ Your Team Name ]", 2.7, 5.35, 3.5, 0.35, font_size=13, bold=True, color=WHITE)
    _add_text(s, "Stock Ticker:", 7.0, 5.1, 3.5, 0.35, font_size=11, color=GRAY)
    _add_text(s, "[ GSE: TICKER ]", 7.0, 5.35, 3.5, 0.35, font_size=13, bold=True, color=GOLD)
    _add_text(s, "Recommendation:", 2.7, 5.8, 3.5, 0.35, font_size=11, color=GRAY)
    _add_text(s, "[ BUY / SELL / HOLD ]", 2.7, 6.05, 3.5, 0.35, font_size=13, bold=True, color=GREEN)
    _add_text(s, "Date:", 7.0, 5.8, 3.5, 0.35, font_size=11, color=GRAY)
    _add_text(s, "[ Presentation Date ]", 7.0, 6.05, 3.5, 0.35, font_size=13, color=WHITE)
    _footer(s, 1)

    # ── SLIDE 2: Table of Contents ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Table of Contents", "Your 10-Section Roadmap", BLUE)
    cols = [
        [("01", "Opening Statement",    "15-second pitch + thesis"),
         ("02", "Economic Analysis",    "Ghana macro environment"),
         ("03", "Industry Analysis",    "Sector & competitive landscape"),
         ("04", "Company Overview",     "Business model & management"),
         ("05", "Financial Summary",    "Key metrics & trends")],
        [("06", "SWOT Analysis",        "Strengths, weaknesses, opportunities, threats"),
         ("07", "Valuation — Model 1",  "DCF / Intrinsic Value"),
         ("08", "Valuation — Model 2",  "DDM / Target Price"),
         ("09", "Valuation Summary",    "Blended target price"),
         ("10", "Recommendation",       "BUY / SELL / HOLD + Conclusion")],
    ]
    for ci, col in enumerate(cols):
        x = 0.3 + ci * 6.55
        for ri, (num, title, desc) in enumerate(col):
            y = 1.5 + ri * 1.05
            _add_rect(s, x, y, 6.2, 0.95, RGBColor(0x1e, 0x29, 0x3b))
            _add_rect(s, x, y, 0.55, 0.95, BLUE)
            _add_text(s, num, x, y + 0.22, 0.55, 0.5,
                      font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _add_text(s, title, x + 0.65, y + 0.05, 5.4, 0.4,
                      font_size=12, bold=True, color=WHITE)
            _add_text(s, desc, x + 0.65, y + 0.5, 5.4, 0.35,
                      font_size=9, color=GRAY)
    _footer(s, 2)

    # ── SLIDE 3: Opening Statement ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Opening Statement", "The 15-Second Pitch — Lead with conviction", GOLD)
    _add_rect(s, 0.4, 1.5, 12.5, 2.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 0.4, 1.5, 0.1, 2.2, GOLD)
    _add_text(s, '"I think that', 0.6, 1.6, 11.8, 0.45, font_size=15, color=GRAY, italic=True)
    _add_text(s, '[ COMPANY NAME ]', 0.6, 1.95, 4.5, 0.5, font_size=20, bold=True, color=GOLD)
    _add_text(s, 'is a', 5.0, 1.95, 1.0, 0.5, font_size=16, color=WHITE, italic=True)
    _add_text(s, '[ BUY / SELL / HOLD ]', 5.9, 1.95, 4.5, 0.5, font_size=20, bold=True, color=GREEN)
    _add_text(s, 'because [ qualitative reason — e.g. dominant market position, strong earnings growth ]', 0.6, 2.5, 11.8, 0.5, font_size=13, color=WHITE, italic=True)
    _add_text(s, 'and our valuation shows the stock is [ undervalued / overvalued ] at its current price of GHc [ X.XX ]."', 0.6, 3.0, 11.8, 0.55, font_size=13, color=LIGHTBLUE, italic=True)
    # Tips
    _add_text(s, "What Makes a Great Opening", 0.4, 3.85, 12, 0.4, font_size=13, bold=True, color=WHITE)
    tips = [
        ("Qualitative", "Explain WHY the company is compelling — moat, brand, leadership, growth story", BLUE),
        ("Quantitative", "Back it with numbers — how much is it undervalued? What's the target price?", GREEN),
        ("Conviction", "Sound confident. A pitch that starts with 'I think' ends with 'I'm sure'", GOLD),
    ]
    for i, (label, tip, col) in enumerate(tips):
        x = 0.4 + i * 4.3
        _add_rect(s, x, 4.3, 4.1, 1.3, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 4.3, 4.1, 0.35, col)
        _add_text(s, label, x + 0.15, 4.32, 3.8, 0.3, font_size=11, bold=True, color=WHITE)
        _add_text(s, tip, x + 0.15, 4.72, 3.8, 0.75, font_size=9.5, color=RGBColor(0xba, 0xd0, 0xf8))
    _footer(s, 3)

    # ── SLIDE 4: Economic Analysis Overview ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Economic Analysis", "Step 1 of 6 — The Macro Environment", BLUE)
    _add_text(s, "Understand the Ghana economic environment your stock operates in.", 0.4, 1.5, 12.5, 0.4, font_size=12, color=GRAY)
    indicators = [
        ("GDP Growth Rate", "[ Enter % ]", "Is the economy expanding?", BLUE),
        ("Inflation Rate (CPI)", "[ Enter % ]", "Affects costs & purchasing power", RED),
        ("BoG Policy Rate", "[ Enter % ]", "Key interest rate driver", GOLD),
        ("GHc/USD Exchange Rate", "[ Enter rate ]", "FX risk for import-heavy firms", GREEN),
        ("Government Debt/GDP", "[ Enter % ]", "Fiscal health indicator", LIGHTBLUE),
        ("Consumer Confidence", "[ High / Med / Low ]", "Household spending outlook", RGBColor(0xa8, 0x5c, 0xf5)),
    ]
    for i, (name, val, desc, col) in enumerate(indicators):
        row, col_i = divmod(i, 3)
        x = 0.4 + col_i * 4.3
        y = 2.0 + row * 2.0
        _add_rect(s, x, y, 4.1, 1.75, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, y, 4.1, 0.32, col)
        _add_text(s, name, x + 0.15, y + 0.04, 3.8, 0.28, font_size=9.5, bold=True, color=WHITE)
        _add_text(s, val, x + 0.15, y + 0.42, 3.8, 0.6, font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_text(s, desc, x + 0.15, y + 1.35, 3.8, 0.3, font_size=9, color=GRAY)
    _footer(s, 4)

    # ── SLIDE 5: Economic Analysis — Impact on Your Stock ───────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Economic Analysis", "How the Macro Environment Affects [ Company Name ]", BLUE)
    _add_text(s, "Map each macro factor to a direct impact on your chosen company.", 0.4, 1.5, 12.5, 0.35, font_size=11, color=GRAY)
    headers = ["Macro Factor", "Current Status", "Impact on Company", "Bullish / Bearish?"]
    widths  = [2.8, 2.5, 5.2, 2.3]
    x_starts = [0.2, 3.05, 5.6, 10.85]
    _add_rect(s, 0.2, 1.95, 13.0, 0.4, BLUE)
    for h, xs, w in zip(headers, x_starts, widths):
        _add_text(s, h, xs + 0.05, 1.98, w - 0.1, 0.34,
                  font_size=10, bold=True, color=WHITE)
    rows = [
        ("Inflation", "[ X% ]", "[ Effect on revenues / costs ]", "[ Bullish ▲ ]"),
        ("Interest Rates", "[ X% ]", "[ Effect on borrowing / capex ]", "[ Bearish ▼ ]"),
        ("Exchange Rate", "[ GHc/USD ]", "[ FX impact on imports/exports ]", "[ Neutral — ]"),
        ("GDP Growth", "[ X% ]", "[ Demand growth for products ]", "[ Bullish ▲ ]"),
        ("Government Policy", "[ Describe ]", "[ Regulation / subsidy impact ]", "[ Describe ]"),
    ]
    for ri, row in enumerate(rows):
        bg = RGBColor(0x1e, 0x29, 0x3b) if ri % 2 == 0 else RGBColor(0x0f, 0x17, 0x2a)
        y = 2.4 + ri * 0.82
        _add_rect(s, 0.2, y, 13.0, 0.78, bg)
        for ci, (cell, xs, w) in enumerate(zip(row, x_starts, widths)):
            c = GOLD if ci == 3 else WHITE
            _add_text(s, cell, xs + 0.05, y + 0.22, w - 0.1, 0.38,
                      font_size=9.5, color=c)
    _footer(s, 5)

    # ── SLIDE 6: Industry Analysis Overview ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Industry Analysis", "Step 2 of 6 — The Competitive Landscape", RGBColor(0x5b, 0x21, 0xb6))
    _add_text(s, "Sector:  [ Banking / Manufacturing / Telecoms / Oil & Gas / Food & Beverage / Other ]", 0.4, 1.5, 12.5, 0.35, font_size=12, bold=True, color=GOLD)
    # Industry metrics
    metrics = [
        ("Industry Size (GHc)", "[ Total market size ]"),
        ("Market Growth Rate", "[ % per year ]"),
        ("Your Company Market Share", "[ X% of sector ]"),
        ("# of Key Competitors", "[ List 3 ]"),
        ("Profit Margin (Industry Avg)", "[ X% ]"),
        ("Regulatory Body", "[ SEC / BoG / NCA / Other ]"),
    ]
    for i, (label, val) in enumerate(metrics):
        row, ci = divmod(i, 2)
        x = 0.4 + ci * 6.4
        y = 2.0 + row * 1.2
        _add_rect(s, x, y, 6.1, 1.05, RGBColor(0x1e, 0x29, 0x3b))
        _add_text(s, label, x + 0.15, y + 0.08, 5.8, 0.35, font_size=10, color=GRAY)
        _add_text(s, val,   x + 0.15, y + 0.45, 5.8, 0.45, font_size=16, bold=True, color=WHITE)
    _footer(s, 6)

    # ── SLIDE 7: Porter's 5 Forces ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Industry Analysis — Porter's 5 Forces",
                "Rate each force: Low / Medium / High", RGBColor(0x5b, 0x21, 0xb6))
    forces = [
        ("Competitive Rivalry",  "How intense is competition among existing players in this sector?",
         "[ Low / Medium / High ]", 0.4, 1.6),
        ("Threat of New Entrants","How easy is it for new companies to enter and compete?",
         "[ Low / Medium / High ]", 7.0, 1.6),
        ("Threat of Substitutes", "Can customers easily switch to an alternative product or service?",
         "[ Low / Medium / High ]", 0.4, 3.8),
        ("Bargaining Power — Suppliers", "Do suppliers have power to raise input prices?",
         "[ Low / Medium / High ]", 7.0, 3.8),
        ("Bargaining Power — Buyers", "Do customers have power to demand lower prices?",
         "[ Low / Medium / High ]", 3.7, 5.8),
    ]
    colors_f = [BLUE, RGBColor(0x5b, 0x21, 0xb6), GREEN, RED, GOLD]
    for (title, desc, rating, x, y), col in zip(forces, colors_f):
        _add_rect(s, x, y, 5.8, 1.6, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, y, 5.8, 0.3, col)
        _add_text(s, title, x + 0.15, y + 0.03, 5.5, 0.28, font_size=10, bold=True, color=WHITE)
        _add_text(s, desc,  x + 0.15, y + 0.38, 5.5, 0.55, font_size=9.5, color=GRAY)
        _add_text(s, rating, x + 0.15, y + 1.1, 5.5, 0.38, font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    _footer(s, 7)

    # ── SLIDE 8: Company Overview ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Company Overview", "Step 3 of 6 — Know Your Company Inside Out", GREEN)
    _add_text(s, "Company:  [ Full Legal Name ]      GSE Ticker:  [ XXX ]      Sector:  [ Sector ]", 0.4, 1.5, 12.5, 0.4, font_size=12, bold=True, color=GOLD)
    info_items = [
        ("Founded", "[ Year ]"),
        ("HQ Location", "[ City, Ghana ]"),
        ("CEO", "[ Name ]"),
        ("Employees", "[ Number ]"),
        ("Listing Date", "[ GSE listing year ]"),
        ("GSE Ticker", "[ XXX ]"),
    ]
    _add_rect(s, 0.4, 2.05, 5.2, 4.8, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Company Facts", 0.55, 2.1, 4.8, 0.4, font_size=12, bold=True, color=LIGHTBLUE)
    for i, (label, val) in enumerate(info_items):
        y = 2.6 + i * 0.7
        _add_text(s, label, 0.6, y, 2.0, 0.35, font_size=10, color=GRAY)
        _add_text(s, val,   2.6, y, 2.8, 0.35, font_size=10, bold=True, color=WHITE)
    _add_rect(s, 5.8, 2.05, 7.1, 2.25, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Business Description", 6.0, 2.12, 6.8, 0.4, font_size=12, bold=True, color=LIGHTBLUE)
    _placeholder_box(s, "Describe the company's core business, products/services, and revenue model in 3-5 sentences", 5.85, 2.55, 6.9, 1.6)
    _add_rect(s, 5.8, 4.45, 7.1, 2.4, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Competitive Advantage / Moat", 6.0, 4.52, 6.8, 0.4, font_size=12, bold=True, color=LIGHTBLUE)
    _placeholder_box(s, "What makes this company hard to compete with? Brand, patents, scale, network effects, cost advantage?", 5.85, 4.95, 6.9, 1.75)
    _footer(s, 8)

    # ── SLIDE 9: Management & Governance ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Management & Corporate Governance", "Who leads the company — and can you trust them?", GREEN)
    _add_text(s, "Strong management is often the deciding factor between a good company and a great investment.", 0.4, 1.5, 12.5, 0.35, font_size=11, color=GRAY, italic=True)
    mgmt = [
        ("Chief Executive Officer (CEO)", "[ Name ]", "[ Key achievements, tenure, track record ]"),
        ("Chief Financial Officer (CFO)", "[ Name ]", "[ Financial expertise, previous roles ]"),
        ("Board Chairman",               "[ Name ]", "[ Independent? Qualifications? ]"),
        ("Largest Shareholder",          "[ Name / Institution ]", "[ % held — alignment of interests? ]"),
    ]
    for i, (role, name, detail) in enumerate(mgmt):
        y = 2.0 + i * 1.25
        _add_rect(s, 0.4, y, 12.5, 1.1, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, 0.4, y, 0.08, 1.1, GREEN)
        _add_text(s, role,   0.6, y + 0.05, 4.0, 0.4, font_size=10, color=GRAY)
        _add_text(s, name,   0.6, y + 0.45, 4.0, 0.5, font_size=14, bold=True, color=WHITE)
        _add_text(s, detail, 4.6, y + 0.25, 8.2, 0.6, font_size=10, color=LIGHTBLUE)
    _footer(s, 9)

    # ── SLIDE 10: Financial Summary ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Financial Summary", "Step 4 of 6 — Key Metrics (from Annual Reports)", GOLD)
    _add_text(s, "Source: Company Annual Reports / GSE Disclosures / Ghana Stock Exchange Data", 0.4, 1.5, 12.5, 0.35, font_size=10, italic=True, color=GRAY)
    # Metrics grid
    kpis = [
        ("Revenue (GHc)", "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", BLUE),
        ("Net Profit (GHc)", "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", GREEN),
        ("EPS (GHc)",        "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", GOLD),
        ("DPS (GHc)",        "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", LIGHTBLUE),
        ("P/E Ratio",        "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", RGBColor(0xa8, 0x5c, 0xf5)),
        ("ROE (%)",          "[ Year-2 ]", "[ Year-1 ]", "[ Latest ]", RED),
    ]
    # Table header
    _add_rect(s, 0.3, 2.0, 12.7, 0.42, NAVY)
    for hdr, xp in [("Metric", 0.4), ("Year -2", 4.2), ("Year -1", 7.2), ("Latest Year", 10.0)]:
        _add_text(s, hdr, xp, 2.02, 3.0, 0.38, font_size=10, bold=True, color=LIGHTBLUE)
    for ri, (name, y2, y1, yl, col) in enumerate(kpis):
        bg = RGBColor(0x1e, 0x29, 0x3b) if ri % 2 == 0 else RGBColor(0x14, 0x1e, 0x2e)
        y = 2.45 + ri * 0.75
        _add_rect(s, 0.3, y, 12.7, 0.7, bg)
        _add_rect(s, 0.3, y, 0.07, 0.7, col)
        _add_text(s, name, 0.5, y + 0.18, 3.5, 0.38, font_size=11, bold=True, color=WHITE)
        for val, xp in [(y2, 4.2), (y1, 7.2), (yl, 10.0)]:
            _add_text(s, val, xp, y + 0.18, 2.8, 0.38, font_size=11, color=col)
    _footer(s, 10)

    # ── SLIDE 11: Revenue & Earnings Trend ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Revenue & Earnings Trend", "5-Year Performance at a Glance", GOLD)
    _add_text(s, "Insert a bar or line chart here showing 5-year revenue and net profit trend.", 0.4, 1.5, 12.5, 0.35, font_size=11, color=GRAY, italic=True)
    _placeholder_box(s, "INSERT CHART: 5-Year Revenue vs Net Profit Bar Chart\n(Use data from Annual Reports — you can embed a chart from Excel or paste an image)", 0.4, 1.95, 8.5, 4.5)
    # Side commentary
    _add_rect(s, 9.1, 1.95, 4.1, 4.5, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Key Observations", 9.25, 2.05, 3.8, 0.4, font_size=12, bold=True, color=LIGHTBLUE)
    obs = ["Revenue CAGR: [ X% over 5 yrs ]",
           "Profit margin trend: [ improving / declining ]",
           "Best year: [ Year ] — why?",
           "Worst year: [ Year ] — why?",
           "Guidance for next year: [ management forecast ]"]
    for i, o in enumerate(obs):
        _add_text(s, f"• {o}", 9.25, 2.55 + i * 0.7, 3.8, 0.6, font_size=9.5, color=WHITE)
    _footer(s, 11)

    # ── SLIDE 12: Ratios & Balance Sheet ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Key Ratios & Balance Sheet", "Profitability, Efficiency, and Leverage", GOLD)
    sections = [
        ("Profitability Ratios", BLUE, [
            ("Gross Margin",    "[ X% ]", "Revenue - COGS / Revenue"),
            ("Net Profit Margin","[ X% ]","Net Profit / Revenue"),
            ("ROE",             "[ X% ]", "Net Profit / Shareholder Equity"),
            ("ROA",             "[ X% ]", "Net Profit / Total Assets"),
        ]),
        ("Valuation Ratios", GOLD, [
            ("P/E Ratio",       "[ Xx ]", "Price / EPS"),
            ("P/B Ratio",       "[ Xx ]", "Price / Book Value per Share"),
            ("EV/EBITDA",       "[ Xx ]", "Enterprise Value / EBITDA"),
            ("Dividend Yield",  "[ X% ]", "DPS / Share Price"),
        ]),
        ("Leverage Ratios", RED, [
            ("Debt-to-Equity",  "[ Xx ]", "Total Debt / Equity"),
            ("Current Ratio",   "[ Xx ]", "Current Assets / Current Liabilities"),
            ("Interest Cover",  "[ Xx ]", "EBIT / Interest Expense"),
            ("Net Debt (GHc)",  "[ Xm ]", "Total Debt - Cash & Equivalents"),
        ]),
    ]
    for ci, (title, col, items) in enumerate(sections):
        x = 0.3 + ci * 4.35
        _add_rect(s, x, 1.5, 4.15, 5.5, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 1.5, 4.15, 0.38, col)
        _add_text(s, title, x + 0.15, 1.53, 3.9, 0.32, font_size=11, bold=True, color=WHITE)
        for ri, (label, val, formula) in enumerate(items):
            y = 2.0 + ri * 1.22
            _add_text(s, label,   x + 0.15, y,        3.8, 0.35, font_size=10, color=GRAY)
            _add_text(s, val,     x + 0.15, y + 0.35, 3.8, 0.45, font_size=18, bold=True, color=col)
            _add_text(s, formula, x + 0.15, y + 0.78, 3.8, 0.3,  font_size=8,  color=RGBColor(0x47, 0x55, 0x69), italic=True)
    _footer(s, 12)

    # ── SLIDE 13: SWOT Analysis ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "SWOT Analysis", "Step 5 — Internal Strengths & Weaknesses, External Opportunities & Threats", GREEN)
    quads = [
        ("S  STRENGTHS",     GREEN, 0.3,  1.45, [
            "[ Strong brand / dominant market position ]",
            "[ Consistent dividend payment history ]",
            "[ Experienced management team ]",
            "[ Low-cost producer / economies of scale ]"]),
        ("W  WEAKNESSES",    RED,   6.85, 1.45, [
            "[ High debt levels ]",
            "[ Heavy reliance on single product/market ]",
            "[ Weak corporate governance ]",
            "[ Thin profit margins vs peers ]"]),
        ("O  OPPORTUNITIES", BLUE,  0.3,  4.45, [
            "[ Economic growth driving demand ]",
            "[ Regulatory changes favouring the sector ]",
            "[ Geographic expansion within Africa ]",
            "[ New product line in development ]"]),
        ("T  THREATS",       GOLD,  6.85, 4.45, [
            "[ New foreign competitors entering Ghana ]",
            "[ Currency depreciation (GHc/USD) ]",
            "[ Rising interest rates increasing debt costs ]",
            "[ Regulatory tightening by SEC / BoG ]"]),
    ]
    for (title, col, x, y, bullets) in quads:
        _add_rect(s, x, y, 6.35, 2.8, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, y, 6.35, 0.38, col)
        _add_text(s, title, x + 0.15, y + 0.05, 6.1, 0.3, font_size=12, bold=True, color=WHITE)
        for i, b in enumerate(bullets):
            _add_text(s, f"•  {b}", x + 0.15, y + 0.5 + i * 0.54, 6.1, 0.48, font_size=10, color=WHITE)
    _footer(s, 13)

    # ── SLIDE 14: Investment Thesis ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Investment Thesis", "Why This Stock? Your 3 Core Arguments", BLUE)
    _add_text(s, "State the 3 most compelling reasons for your recommendation. Each argument must combine qualitative insight with data.", 0.4, 1.5, 12.5, 0.4, font_size=11, color=GRAY, italic=True)
    args = [
        ("Argument 1 — [ Growth ]",
         "[ Describe first reason — e.g. 'Revenue growing at 25% CAGR driven by digital banking adoption' ]",
         "Supporting data: [ Specific metric ]", BLUE),
        ("Argument 2 — [ Valuation ]",
         "[ Describe second reason — e.g. 'Stock trades at 8x P/E vs sector average of 14x, offering 40% discount' ]",
         "Supporting data: [ Specific metric ]", GOLD),
        ("Argument 3 — [ Catalyst ]",
         "[ Describe third reason — e.g. 'Upcoming bond issuance and FX recovery will expand margins by 3pp' ]",
         "Supporting data: [ Specific metric ]", GREEN),
    ]
    for i, (title, body, data, col) in enumerate(args):
        y = 2.05 + i * 1.55
        _add_rect(s, 0.4, y, 12.5, 1.4, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, 0.4, y, 0.1,  1.4, col)
        _add_text(s, title, 0.65, y + 0.05, 11.8, 0.4, font_size=13, bold=True, color=col)
        _add_text(s, body,  0.65, y + 0.47, 11.8, 0.5, font_size=11, color=WHITE)
        _add_text(s, data,  0.65, y + 0.97, 11.8, 0.35, font_size=10, italic=True, color=GRAY)
    _footer(s, 14)

    # ── SLIDE 15: Valuation Introduction ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Valuation Framework", "Step 6 — How We Value the Stock", RGBColor(0xf9, 0x73, 0x16))
    _add_text(s, "We used 3 complementary models. Each provides a different lens. Together they give a robust target price range.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    models = [
        ("DCF Valuation", "Discounted Cash Flow", "Intrinsic value based on projected future free cash flows discounted at WACC.", "Best for: companies with predictable FCF", BLUE, "Slide 16"),
        ("Graham Formula", "Intrinsic Value (Benjamin Graham)", "Classic margin-of-safety formula: √(22.5 × EPS × BVPS)", "Best for: any listed company with positive earnings", GREEN, "Slide 17"),
        ("Target Price / DDM", "Forward EPS × P/E or Gordon Growth", "Analyst-style price target or dividend discount model.", "Best for: dividend-paying stocks or peer multiples", GOLD, "Slide 18"),
    ]
    for i, (short, full, desc, note, col, sref) in enumerate(models):
        x = 0.4 + i * 4.3
        _add_rect(s, x, 2.05, 4.1, 4.6, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 2.05, 4.1, 0.35, col)
        _add_text(s, sref,  x + 3.0, 2.07, 1.0, 0.28, font_size=8, color=WHITE)
        _add_text(s, short, x + 0.15, 2.07, 3.8, 0.28, font_size=11, bold=True, color=WHITE)
        _add_text(s, full,  x + 0.15, 2.5,  3.8, 0.5,  font_size=10, bold=True, color=col)
        _add_text(s, desc,  x + 0.15, 3.08, 3.8, 1.2,  font_size=9.5, color=WHITE)
        _add_text(s, note,  x + 0.15, 4.35, 3.8, 0.5,  font_size=8.5, italic=True, color=GRAY)
        _add_text(s, "Target Price", x + 0.15, 5.2, 3.8, 0.3, font_size=9, color=GRAY)
        _add_text(s, "GHc [ X.XX ]", x + 0.15, 5.5, 3.8, 0.45, font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    _footer(s, 15)

    # ── SLIDE 16: DCF Valuation ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Valuation Model 1 — DCF", "Discounted Cash Flow Analysis", BLUE)
    _add_text(s, "The DCF model values the company based on its ability to generate free cash flow in the future, discounted to today's value.", 0.4, 1.5, 12.5, 0.35, font_size=10.5, italic=True, color=GRAY)
    # Inputs
    _add_rect(s, 0.4, 2.0, 5.5, 5.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 0.4, 2.0, 5.5, 0.38, BLUE)
    _add_text(s, "DCF Inputs", 0.6, 2.03, 5.2, 0.32, font_size=12, bold=True, color=WHITE)
    dcf_inputs = [
        ("Base-Year Free Cash Flow (GHc)", "[ Enter FCF ]"),
        ("Revenue Growth Rate — Yr 1-5",   "[ X% ]"),
        ("Revenue Growth Rate — Yr 6-10",  "[ X% ]"),
        ("EBITDA Margin",                   "[ X% ]"),
        ("Tax Rate",                        "[ X% ]"),
        ("Capex as % of Revenue",           "[ X% ]"),
        ("WACC (Discount Rate)",            "[ X% ]"),
        ("Terminal Growth Rate",            "[ X% ]"),
        ("Shares Outstanding",              "[ X million ]"),
    ]
    for i, (label, val) in enumerate(dcf_inputs):
        y = 2.48 + i * 0.52
        _add_text(s, label, 0.6,  y, 3.5, 0.4, font_size=9.5, color=GRAY)
        _add_text(s, val,   4.0,  y, 1.7, 0.4, font_size=10, bold=True, color=LIGHTBLUE)
    # Outputs
    _add_rect(s, 6.1, 2.0, 6.8, 5.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 6.1, 2.0, 6.8, 0.38, BLUE)
    _add_text(s, "DCF Outputs", 6.3, 2.03, 6.5, 0.32, font_size=12, bold=True, color=WHITE)
    dcf_out = [
        ("PV of FCF (Years 1–10)", "[ GHc Xm ]", LIGHTBLUE),
        ("Terminal Value",          "[ GHc Xm ]", LIGHTBLUE),
        ("Enterprise Value",        "[ GHc Xm ]", WHITE),
        ("Less: Net Debt",          "[ (GHc Xm) ]", RED),
        ("Equity Value",            "[ GHc Xm ]", GREEN),
        ("÷ Shares Outstanding",    "[ X million ]", GRAY),
        ("",                        "",              WHITE),
        ("DCF Intrinsic Value / Share", "GHc [ X.XX ]", GOLD),
        ("Current Market Price",    "GHc [ X.XX ]", WHITE),
        ("Upside / (Downside)",     "[ +X% / -X% ]", GREEN),
    ]
    for i, (label, val, col) in enumerate(dcf_out):
        y = 2.48 + i * 0.46
        if i == 7:
            _add_rect(s, 6.15, y - 0.05, 6.7, 0.56, RGBColor(0x1a, 0x2a, 0x10))
        _add_text(s, label, 6.3, y, 4.8, 0.4, font_size=10 if i != 7 else 12, bold=(i==7), color=GRAY if i not in [7,8,9] else col)
        _add_text(s, val,   10.6, y, 2.2, 0.4, font_size=10 if i != 7 else 15, bold=(i in [7,8,9]), color=col, align=PP_ALIGN.RIGHT)
    _add_text(s, "⚠  Use InvestIQ DCF Calculator at investright.onrender.com", 0.4, 7.0, 12.5, 0.3, font_size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    _footer(s, 16)

    # ── SLIDE 17: Graham / Intrinsic Value ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Valuation Model 2 — Benjamin Graham Formula", "Intrinsic Value & Margin of Safety", GREEN)
    _add_rect(s, 0.4, 1.55, 12.5, 1.0, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Graham Formula:  V = EPS × (8.5 + 2g)  ×  (4.4 / AAA Bond Yield)", 0.6, 1.6, 12, 0.4, font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _add_text(s, "where  g = expected EPS growth rate (%),  8.5 = base P/E for zero-growth company", 0.6, 2.0, 12, 0.35, font_size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    inputs_g = [
        ("Earnings Per Share (EPS)", "[ GHc X.XX ]"),
        ("Expected EPS Growth Rate (g)", "[ X% per year ]"),
        ("Risk-Free / AAA Bond Yield", "[ X% ]"),
        ("Current Market Price", "[ GHc X.XX ]"),
    ]
    _add_rect(s, 0.4, 2.55, 5.8, 4.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Inputs", 0.6, 2.6, 5.5, 0.38, font_size=12, bold=True, color=LIGHTBLUE)
    for i, (label, val) in enumerate(inputs_g):
        y = 3.1 + i * 0.9
        _add_text(s, label, 0.6, y, 3.8, 0.35, font_size=10, color=GRAY)
        _add_text(s, val,   0.6, y + 0.38, 3.8, 0.42, font_size=15, bold=True, color=WHITE)
    _add_rect(s, 6.4, 2.55, 6.5, 4.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Results", 6.6, 2.6, 6.2, 0.38, font_size=12, bold=True, color=LIGHTBLUE)
    results = [
        ("Graham Intrinsic Value", "GHc [ X.XX ]", GREEN),
        ("Current Market Price",   "GHc [ X.XX ]", WHITE),
        ("Margin of Safety",       "[ +X% Discount / -X% Premium ]", GOLD),
        ("Verdict",                "[ Undervalued → BUY  /  Overvalued → SELL ]", LIGHTBLUE),
    ]
    for i, (label, val, col) in enumerate(results):
        y = 3.1 + i * 0.97
        _add_text(s, label, 6.6, y, 6.0, 0.35, font_size=10, color=GRAY)
        _add_text(s, val, 6.6, y + 0.38, 6.0, 0.48, font_size=14, bold=True, color=col)
    _add_text(s, "⚠  Use InvestIQ Intrinsic Value Calculator at investright.onrender.com", 0.4, 6.95, 12.5, 0.3, font_size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    _footer(s, 17)

    # ── SLIDE 18: DDM / Target Price ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Valuation Model 3 — Target Price & DDM", "Forward P/E Target Price + Dividend Discount Model", GOLD)
    _add_rect(s, 0.4, 1.55, 6.1, 5.65, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 0.4, 1.55, 6.1, 0.35, GOLD)
    _add_text(s, "Target Price — Forward P/E Method", 0.6, 1.57, 5.8, 0.3, font_size=11, bold=True, color=WHITE)
    _add_text(s, "Formula:  Target Price = Forward EPS × Target P/E", 0.6, 2.0, 5.8, 0.38, font_size=10, bold=True, color=GOLD)
    tp_inputs = [
        ("Current EPS (GHc)", "[ X.XX ]"),
        ("EPS Growth Rate",   "[ X% ]"),
        ("Forward EPS (1yr)", "[ GHc X.XX ]"),
        ("Sector Average P/E","[ Xx ]"),
        ("Target P/E Applied","[ Xx ]"),
        ("Target Price",      "GHc [ X.XX ]"),
        ("Current Price",     "GHc [ X.XX ]"),
        ("Upside Potential",  "[ +X% ]"),
    ]
    for i, (label, val) in enumerate(tp_inputs):
        y = 2.5 + i * 0.6
        big = i in [5, 7]
        _add_text(s, label, 0.6, y, 3.5, 0.45, font_size=10, color=GRAY)
        _add_text(s, val,   4.0, y, 2.3, 0.45, font_size=13 if big else 10, bold=big, color=GOLD if big else WHITE)
    # DDM side
    _add_rect(s, 6.7, 1.55, 6.3, 5.65, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 6.7, 1.55, 6.3, 0.35, LIGHTBLUE)
    _add_text(s, "Dividend Discount Model (Gordon Growth)", 6.9, 1.57, 6.0, 0.3, font_size=11, bold=True, color=WHITE)
    _add_text(s, "Formula:  P = D₁ / (Ke − g)  where  D₁ = DPS × (1+g)", 6.9, 2.0, 6.0, 0.38, font_size=10, bold=True, color=LIGHTBLUE)
    ddm_inputs = [
        ("Current DPS (GHc)",       "[ X.XX ]"),
        ("Dividend Growth Rate (g)", "[ X% ]"),
        ("D₁ — Next Year DPS",      "GHc [ X.XX ]"),
        ("Cost of Equity (Ke)",      "[ X% ]  (from CAPM)"),
        ("DDM Intrinsic Value",      "GHc [ X.XX ]"),
        ("Current Market Price",     "GHc [ X.XX ]"),
        ("Upside / Downside",        "[ +X% / -X% ]"),
        ("Dividend Yield",           "[ X% ]"),
    ]
    for i, (label, val) in enumerate(ddm_inputs):
        y = 2.5 + i * 0.6
        big = i in [4, 6]
        _add_text(s, label, 6.9, y, 4.0, 0.45, font_size=10, color=GRAY)
        _add_text(s, val,   10.5, y, 2.3, 0.45, font_size=13 if big else 10, bold=big, color=LIGHTBLUE if big else WHITE)
    _add_text(s, "⚠  Use InvestIQ Target Price & DDM Calculators at investright.onrender.com", 0.4, 7.03, 12.5, 0.3, font_size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    _footer(s, 18)

    # ── SLIDE 19: Valuation Summary ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Valuation Summary", "Blended Target Price — Three Models, One Verdict", RGBColor(0xf9, 0x73, 0x16))
    _add_text(s, "We ran three independent valuation models. The blended average gives our final 12-month price target.", 0.4, 1.55, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    # Three model summary boxes
    for i, (model, val, col) in enumerate([
        ("DCF Valuation",            "GHc [ X.XX ]", BLUE),
        ("Graham Intrinsic Value",   "GHc [ X.XX ]", GREEN),
        ("Target Price / DDM",       "GHc [ X.XX ]", GOLD),
    ]):
        x = 0.4 + i * 4.3
        _add_rect(s, x, 2.05, 4.1, 2.0, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 2.05, 4.1, 0.35, col)
        _add_text(s, model, x + 0.15, 2.07, 3.9, 0.3, font_size=10, bold=True, color=WHITE)
        _add_text(s, val, x + 0.15, 2.55, 3.9, 0.85, font_size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_text(s, "[ Upside: +X% ]", x + 0.15, 3.5, 3.9, 0.35, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
    # Blended target box
    _add_rect(s, 2.3, 4.25, 8.7, 2.0, RGBColor(0x14, 0x28, 0x14))
    _add_rect(s, 2.3, 4.25, 8.7, 0.06, GREEN)
    _add_rect(s, 2.3, 6.19, 8.7, 0.06, GREEN)
    _add_text(s, "BLENDED TARGET PRICE (12-Month)", 2.5, 4.3, 8.4, 0.45, font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _add_text(s, "GHc [ X.XX ]", 2.5, 4.85, 8.4, 0.85, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "vs. Current Price GHc [ X.XX ]  ·  Implied Upside: [ +X% ]  ·  12-Month Horizon", 2.5, 5.8, 8.4, 0.35, font_size=10, color=GREEN, align=PP_ALIGN.CENTER)
    _footer(s, 19)

    # ── SLIDE 20: Comparable Companies ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Comparable Companies Analysis", "How Does Our Stock Stack Up Against Peers?", BLUE)
    _add_text(s, "Compare your stock against 3–4 sector peers listed on the GSE on key valuation and profitability metrics.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    headers_c = ["Company", "Ticker", "Price (GHc)", "P/E", "P/B", "ROE (%)", "Div Yield (%)", "Rec."]
    widths_c  = [3.0, 1.3, 1.8, 1.3, 1.3, 1.5, 2.0, 1.5]
    x_pos = [0.2]
    for w in widths_c[:-1]:
        x_pos.append(x_pos[-1] + w)
    _add_rect(s, 0.2, 2.0, 13.0, 0.42, BLUE)
    for hdr, xp, w in zip(headers_c, x_pos, widths_c):
        _add_text(s, hdr, xp + 0.05, 2.02, w - 0.1, 0.36, font_size=10, bold=True, color=WHITE)
    rows_c = [
        ("[ Your Company ]", "[ XXX ]", "[ X.XX ]", "[ Xx ]", "[ Xx ]", "[ X% ]", "[ X% ]", "★ BUY"),
        ("[ Peer 1 ]",       "[ XXX ]", "[ X.XX ]", "[ Xx ]", "[ Xx ]", "[ X% ]", "[ X% ]", "Hold"),
        ("[ Peer 2 ]",       "[ XXX ]", "[ X.XX ]", "[ Xx ]", "[ Xx ]", "[ X% ]", "[ X% ]", "Hold"),
        ("[ Peer 3 ]",       "[ XXX ]", "[ X.XX ]", "[ Xx ]", "[ Xx ]", "[ X% ]", "[ X% ]", "Sell"),
        ("Sector Average",   "—",       "—",         "[ Xx ]", "[ Xx ]", "[ X% ]", "[ X% ]", "—"),
    ]
    for ri, row in enumerate(rows_c):
        bg = RGBColor(0x0a, 0x20, 0x10) if ri == 0 else (RGBColor(0x1e, 0x29, 0x3b) if ri % 2 else RGBColor(0x14, 0x1e, 0x2e))
        y = 2.48 + ri * 0.82
        _add_rect(s, 0.2, y, 13.0, 0.76, bg)
        if ri == 0:
            _add_rect(s, 0.2, y, 0.07, 0.76, GREEN)
        for ci, (cell, xp, w) in enumerate(zip(row, x_pos, widths_c)):
            fc = GREEN if (ri == 0) else (GOLD if ci == 7 and "BUY" in cell else WHITE)
            _add_text(s, cell, xp + 0.05, y + 0.2, w - 0.1, 0.38, font_size=10, bold=(ri == 0), color=fc)
    _footer(s, 20)

    # ── SLIDE 21: Risk Analysis ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Risk Analysis", "Key Risks to Our Investment Thesis", RED)
    _add_text(s, "Every investment has risks. Acknowledging and quantifying them shows the judges you've done thorough work.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    risks = [
        ("FX / Currency Risk",     "HIGH",   RED,   "GHc depreciation raises import costs and erodes margins for companies with USD-linked inputs.",    "Hedge: FX forwards / natural hedging via export revenues"),
        ("Interest Rate Risk",      "MEDIUM", GOLD,  "Rising BoG policy rates increase borrowing costs, squeeze margins, and reduce equity valuations.", "Monitor BoG MPC meetings; prefer low-debt companies"),
        ("Regulatory Risk",         "LOW",    GREEN, "SEC or sector regulator (NCA, BoG) could impose capital requirements or pricing limits.",          "Track regulatory pipeline; assess management relationships"),
        ("Competitive Risk",        "MEDIUM", GOLD,  "New domestic or foreign entrants could erode market share and pricing power.",                      "Assess barriers to entry; moat durability"),
        ("Company-Specific Risk",   "HIGH",   RED,   "[ Describe specific risk unique to this company — earnings miss, debt maturity, key-person risk ]", "[ How does management mitigate this? ]"),
    ]
    for i, (name, level, col, desc, mitig) in enumerate(risks):
        y = 2.0 + i * 1.02
        _add_rect(s, 0.3, y, 12.7, 0.92, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, 0.3, y, 0.07, 0.92, col)
        _add_text(s, name,  0.5, y + 0.04, 3.0, 0.38, font_size=11, bold=True, color=WHITE)
        _add_rect(s, 3.6, y + 0.12, 1.4, 0.45, col)
        _add_text(s, level, 3.65, y + 0.15, 1.3, 0.38, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, desc,  5.2, y + 0.04, 5.5, 0.45, font_size=9.5, color=GRAY)
        _add_text(s, f"Mitigation: {mitig}", 5.2, y + 0.52, 7.6, 0.35, font_size=8.5, italic=True, color=GREEN)
    _footer(s, 21)

    # ── SLIDE 22: Catalysts ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Catalysts & Upside Drivers", "What Will Push the Stock to Our Target Price?", GREEN)
    _add_text(s, "Catalysts are events or developments that will unlock value and drive the share price toward your target within 12 months.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    cats = [
        ("Near-Term (0–3 months)",   GOLD,  [
            "[ Upcoming quarterly results — expected earnings beat ]",
            "[ Dividend announcement / special dividend ]",
            "[ Regulatory approval for new product/branch ]"]),
        ("Medium-Term (3–9 months)", BLUE,  [
            "[ New product launch / geographic expansion ]",
            "[ Debt refinancing at lower rates ]",
            "[ Industry tailwind from BoG policy change ]"]),
        ("Long-Term (9–18 months)",  GREEN, [
            "[ Strategic acquisition or partnership ]",
            "[ Market share gains from competitor weakness ]",
            "[ Macroeconomic recovery driving sector re-rating ]"]),
    ]
    for i, (period, col, bullets) in enumerate(cats):
        x = 0.4 + i * 4.3
        _add_rect(s, x, 2.05, 4.1, 4.7, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 2.05, 4.1, 0.38, col)
        _add_text(s, period, x + 0.15, 2.08, 3.8, 0.3, font_size=11, bold=True, color=WHITE)
        for j, b in enumerate(bullets):
            _add_text(s, f"→  {b}", x + 0.15, 2.6 + j * 1.3, 3.85, 1.18, font_size=9.5, color=WHITE)
    _footer(s, 22)

    # ── SLIDE 23: Sensitivity Analysis ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Sensitivity Analysis", "How Does Our Target Price Change Under Different Scenarios?", BLUE)
    _add_text(s, "A sensitivity table shows judges you've stress-tested your valuation. Vary WACC and growth rate to see price range.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    _add_text(s, "DCF Sensitivity: Intrinsic Value (GHc) — Varying WACC vs Terminal Growth Rate", 0.4, 2.0, 12.5, 0.38, font_size=12, bold=True, color=WHITE)
    wacc_vals = ["WACC - 2%", "WACC - 1%", "Base WACC", "WACC + 1%", "WACC + 2%"]
    tgr_vals  = ["TGR 1%", "TGR 2%", "TGR 3%", "TGR 4%", "TGR 5%"]
    col_w = 2.0
    _add_rect(s, 0.3, 2.55, 13.0, 0.45, BLUE)
    _add_text(s, "↓ WACC   /   TGR →", 0.4, 2.57, 2.1, 0.38, font_size=9, bold=True, color=WHITE)
    for ci, tgr in enumerate(tgr_vals):
        _add_text(s, tgr, 2.5 + ci * col_w, 2.57, col_w - 0.1, 0.38, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, wacc in enumerate(wacc_vals):
        bg = RGBColor(0x0a, 0x20, 0x10) if ri == 2 else (RGBColor(0x1e, 0x29, 0x3b) if ri % 2 else RGBColor(0x14, 0x1e, 0x2e))
        y = 3.05 + ri * 0.75
        _add_rect(s, 0.3, y, 13.0, 0.7, bg)
        _add_text(s, wacc, 0.4, y + 0.16, 2.0, 0.38, font_size=9.5, bold=(ri==2), color=GOLD if ri==2 else WHITE)
        for ci in range(5):
            fc = GREEN if ri == 2 and ci == 2 else (LIGHTBLUE if ri < 2 else (RED if ri > 2 else WHITE))
            _add_text(s, "GHc [ X.XX ]", 2.5 + ci * col_w, y + 0.16, col_w - 0.1, 0.38,
                      font_size=9.5, bold=(ri==2 and ci==2), color=fc, align=PP_ALIGN.CENTER)
    _add_text(s, "Green = Base Case  ·  Blue = Bull Scenario  ·  Red = Bear Scenario", 0.4, 6.95, 12.5, 0.3, font_size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    _footer(s, 23)

    # ── SLIDE 24: Dividend & Cash Flow ───────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Dividend & Cash Flow Analysis", "Is the Dividend Safe? Is the Business Self-Funding?", GOLD)
    _add_text(s, "Strong free cash flow and sustainable dividends signal a financially healthy company.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    _add_rect(s, 0.4, 1.95, 6.0, 5.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 0.4, 1.95, 6.0, 0.38, GOLD)
    _add_text(s, "Dividend History", 0.6, 1.97, 5.7, 0.32, font_size=12, bold=True, color=WHITE)
    div_data = [
        ("Year", "DPS (GHc)", "Payout Ratio", "Yield"),
        ("[ Year -4 ]", "[ X.XX ]", "[ X% ]", "[ X% ]"),
        ("[ Year -3 ]", "[ X.XX ]", "[ X% ]", "[ X% ]"),
        ("[ Year -2 ]", "[ X.XX ]", "[ X% ]", "[ X% ]"),
        ("[ Year -1 ]", "[ X.XX ]", "[ X% ]", "[ X% ]"),
        ("[ Latest ]",  "[ X.XX ]", "[ X% ]", "[ X% ]"),
    ]
    for ri, row in enumerate(div_data):
        bg = BLUE if ri == 0 else (RGBColor(0x1a, 0x20, 0x2e) if ri % 2 else RGBColor(0x14, 0x1a, 0x28))
        y = 2.42 + ri * 0.72
        _add_rect(s, 0.45, y, 5.9, 0.67, bg)
        for ci, cell in enumerate(row):
            _add_text(s, cell, 0.55 + ci * 1.45, y + 0.15, 1.35, 0.38,
                      font_size=9.5 if ri == 0 else 10,
                      bold=(ri == 0), color=WHITE if ri == 0 else GOLD if ci == 1 else WHITE)
    _add_rect(s, 6.6, 1.95, 6.4, 5.2, RGBColor(0x1e, 0x29, 0x3b))
    _add_rect(s, 6.6, 1.95, 6.4, 0.38, GREEN)
    _add_text(s, "Free Cash Flow Summary", 6.8, 1.97, 6.1, 0.32, font_size=12, bold=True, color=WHITE)
    fcf_items = [
        ("Operating Cash Flow",   "GHc [ Xm ]"),
        ("Less: Capital Expenditure", "(GHc [ Xm ])"),
        ("Free Cash Flow",        "GHc [ Xm ]"),
        ("FCF per Share",         "GHc [ X.XX ]"),
        ("FCF Yield",             "[ X% ]"),
        ("FCF Payout Ratio",      "[ X% ]"),
        ("Is Dividend Covered?",  "[ YES / NO ]"),
    ]
    for i, (label, val) in enumerate(fcf_items):
        y = 2.5 + i * 0.65
        sep = i == 2
        if sep:
            _add_rect(s, 6.65, y - 0.07, 6.3, 0.04, GREEN)
        _add_text(s, label, 6.8, y, 4.2, 0.42, font_size=10, bold=sep, color=GRAY if not sep else WHITE)
        _add_text(s, val,   10.7, y, 2.1, 0.42, font_size=12 if sep else 10, bold=sep, color=GREEN if sep else WHITE, align=PP_ALIGN.RIGHT)
    _footer(s, 24)

    # ── SLIDE 25: Scenario Analysis ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Scenario Analysis", "Bull, Base & Bear Case — What Are the Outcomes?", BLUE)
    _add_text(s, "Present 3 scenarios. The base case is your primary recommendation. Bull and bear show the range of possible outcomes.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    scenarios = [
        ("BEAR CASE",    RED,   "[ Probability: X% ]",
         ["Revenue declines X% due to macro headwinds",
          "Margin compression from rising costs",
          "Dividend cut or suspension"],
         "GHc [ X.XX ]", "Downside of [ -X% ]"),
        ("BASE CASE",    GOLD,  "[ Probability: X% ]",
         ["Revenue grows at X% — in line with guidance",
          "Margins stable at current levels",
          "Dividend maintained at current DPS"],
         "GHc [ X.XX ]", "Upside of [ +X% ]"),
        ("BULL CASE",    GREEN, "[ Probability: X% ]",
         ["Revenue beats expectations by X%",
          "Margin expansion from operating leverage",
          "Re-rating as sector recovers"],
         "GHc [ X.XX ]", "Upside of [ +X% ]"),
    ]
    for i, (title, col, prob, bullets, tp, updown) in enumerate(scenarios):
        x = 0.4 + i * 4.3
        _add_rect(s, x, 2.05, 4.1, 5.1, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, 2.05, 4.1, 0.38, col)
        _add_text(s, title, x + 0.15, 2.07, 3.8, 0.3, font_size=12, bold=True, color=WHITE)
        _add_text(s, prob,  x + 0.15, 2.5,  3.8, 0.3, font_size=9.5, color=col)
        for j, b in enumerate(bullets):
            _add_text(s, f"• {b}", x + 0.15, 2.92 + j * 0.62, 3.85, 0.55, font_size=9.5, color=WHITE)
        _add_rect(s, x + 0.15, 4.95, 3.8, 1.05, RGBColor(0x0a, 0x10, 0x1a))
        _add_text(s, "Target Price", x + 0.15, 4.98, 3.8, 0.3, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)
        _add_text(s, tp,     x + 0.15, 5.3, 3.8, 0.42, font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_text(s, updown, x + 0.15, 5.75, 3.8, 0.3, font_size=10, color=col, align=PP_ALIGN.CENTER)
    _footer(s, 25)

    # ── SLIDE 26: Recommendation ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Our Recommendation", "The Final Verdict — BUY / SELL / HOLD", GREEN)
    # Large verdict box
    _add_rect(s, 0.4, 1.6, 5.5, 5.5, RGBColor(0x0a, 0x20, 0x10))
    _add_rect(s, 0.4, 1.6, 5.5, 0.06, GREEN)
    _add_rect(s, 0.4, 7.04, 5.5, 0.06, GREEN)
    _add_text(s, "RECOMMENDATION", 0.5, 1.75, 5.3, 0.45, font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _add_text(s, "[ BUY ]", 0.5, 2.3, 5.3, 1.3, font_size=60, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _add_text(s, "[ Company Name ]  ·  GSE: [ XXX ]", 0.5, 3.7, 5.3, 0.45, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    key_metrics = [
        ("Current Price",  "GHc [ X.XX ]"),
        ("Target Price",   "GHc [ X.XX ]"),
        ("Upside",         "[ +X% ]"),
        ("Time Horizon",   "12 months"),
        ("Risk Rating",    "[ Low / Med / High ]"),
    ]
    for i, (label, val) in enumerate(key_metrics):
        y = 4.28 + i * 0.48
        _add_text(s, label, 0.6, y, 2.5, 0.38, font_size=10, color=GRAY)
        _add_text(s, val,   3.0, y, 2.7, 0.38, font_size=11, bold=True, color=GREEN if i == 2 else WHITE, align=PP_ALIGN.RIGHT)
    # Right side — reasoning
    _add_rect(s, 6.1, 1.6, 7.1, 5.5, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "Why We Recommend This", 6.3, 1.72, 6.8, 0.42, font_size=13, bold=True, color=WHITE)
    reasons = [
        ("1. Valuation", "[ Stock trades at significant discount to intrinsic value of GHc X.XX — X% margin of safety ]"),
        ("2. Earnings Growth", "[ EPS growing at X% CAGR with improving margins driven by [ specific driver ] ]"),
        ("3. Strong Dividends", "[ Consistent dividend payer — X% yield, well-covered by FCF at X% payout ratio ]"),
        ("4. Sector Catalyst", "[ [ Industry-specific catalyst ] will drive a re-rating within 12 months ]"),
    ]
    for i, (label, text) in enumerate(reasons):
        y = 2.2 + i * 1.2
        _add_text(s, label, 6.3, y, 6.8, 0.38, font_size=11, bold=True, color=GOLD)
        _add_text(s, text,  6.3, y + 0.4, 6.8, 0.65, font_size=9.5, color=WHITE)
    _footer(s, 26)

    # ── SLIDE 27: Key Risks to Our Thesis ───────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Key Risks to Our Recommendation", "When Would We Change Our View?", RED)
    _add_text(s, "Be honest about what could go wrong. Judges respect intellectual honesty — and want to know you've thought of the downside.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    risks2 = [
        ("Risk 1 — [ Main Risk ]",       "[ Describe the primary risk to your thesis ]",        "[ What event would confirm this risk? ]", "[ What would you do — exit position? ]"),
        ("Risk 2 — [ Secondary Risk ]",  "[ Describe the secondary risk ]",                      "[ Trigger event ]",                       "[ Response strategy ]"),
        ("Risk 3 — [ Company-Specific ]","[ Describe the company-specific risk ]",               "[ What to monitor in quarterly results? ]","[ Stop-loss or thesis review trigger ]"),
    ]
    _add_rect(s, 0.2, 2.0, 12.9, 0.38, RGBColor(0x33, 0x00, 0x00))
    for hdr, xp, w in [("Risk", 0.3, 4.0), ("If This Happens", 4.4, 3.5), ("Monitor This", 8.0, 2.7), ("Our Response", 10.8, 2.2)]:
        _add_text(s, hdr, xp, 2.02, w, 0.34, font_size=10, bold=True, color=WHITE)
    for ri, (risk, impact, monitor, response) in enumerate(risks2):
        y = 2.48 + ri * 1.55
        bg = RGBColor(0x2a, 0x10, 0x10) if ri % 2 == 0 else RGBColor(0x1e, 0x29, 0x3b)
        _add_rect(s, 0.2, y, 12.9, 1.42, bg)
        _add_rect(s, 0.2, y, 0.07, 1.42, RED)
        _add_text(s, risk,     0.35, y + 0.05, 3.9, 0.38, font_size=10, bold=True, color=RED)
        _add_text(s, impact,   0.35, y + 0.5,  3.9, 0.8,  font_size=9.5, color=WHITE)
        _add_text(s, monitor,  4.4,  y + 0.3,  3.4, 0.8,  font_size=9.5, color=GRAY)
        _add_text(s, response, 8.0,  y + 0.3,  4.9, 0.8,  font_size=9.5, color=GRAY)
    _footer(s, 27)

    # ── SLIDE 28: Conclusion ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Conclusion", "Tying It All Together — Final 2 Minutes", BLUE)
    _add_text(s, "Your conclusion restates the thesis, reinforces the 3 strongest arguments, and closes with a clear, confident verdict.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    # Conclusion structure
    steps = [
        ("Restate Your Thesis",    BLUE,  "[ We recommended [Company] as a BUY. Our analysis shows the stock is trading at a X% discount to its intrinsic value of GHc X.XX, offering compelling upside within 12 months. ]"),
        ("Top 3 Arguments",        GOLD,  "[ 1: [Strongest valuation argument] ]  ·  [ 2: [Earnings/growth argument] ]  ·  [ 3: [Catalyst/sector argument] ]"),
        ("Key Risk Acknowledged",  RED,   "[ The primary risk to our thesis is [X]. However, we believe [mitigation factor] limits the downside to [Y%]. ]"),
        ("Final Verdict & Close",  GREEN, "[ Based on our comprehensive analysis, we are confident in our BUY recommendation with a 12-month target price of GHc [X.XX]. Thank you — we welcome your questions. ]"),
    ]
    for i, (title, col, text) in enumerate(steps):
        y = 2.05 + i * 1.28
        _add_rect(s, 0.4, y, 12.5, 1.15, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, 0.4, y, 0.09, 1.15, col)
        num_box_x = 0.6
        _add_rect(s, num_box_x, y + 0.28, 0.5, 0.5, col)
        _add_text(s, str(i+1), num_box_x, y + 0.28, 0.5, 0.5, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, title, 1.25, y + 0.05, 11.4, 0.38, font_size=12, bold=True, color=col)
        _add_text(s, text,  1.25, y + 0.52, 11.4, 0.58, font_size=10, color=WHITE)
    _footer(s, 28)

    # ── SLIDE 29: Appendix ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _header_bar(s, "Appendix", "Supporting Data & Additional Analysis", GRAY)
    _add_text(s, "Include additional charts, tables, or data here. Judges may ask questions from your appendix during Q&A.", 0.4, 1.5, 12.5, 0.35, font_size=11, italic=True, color=GRAY)
    appendix_items = [
        ("A.1", "Full 5-Year Income Statement",    "Detailed P&L showing revenue, COGS, EBIT, net profit for 5 years"),
        ("A.2", "Full Balance Sheet Summary",       "Total assets, liabilities, equity, net debt breakdown"),
        ("A.3", "Detailed DCF Model Workings",      "Full year-by-year FCF projections, terminal value computation"),
        ("A.4", "Cost of Equity Calculation",       "CAPM: Risk-free rate + Beta × (Market Return - Risk-free rate)"),
        ("A.5", "Peer Comparison — Full Table",     "Extended comparable companies analysis with all metrics"),
        ("A.6", "Historical Share Price Chart",     "3–5 year price performance vs GSE Composite Index"),
        ("A.7", "Management Profiles",              "Detailed bios for CEO, CFO, and Board members"),
        ("A.8", "Sources & References",             "Annual reports, GSE data, news sources, research reports cited"),
    ]
    for i, (ref, title, desc) in enumerate(appendix_items):
        row, ci = divmod(i, 2)
        x = 0.4 + ci * 6.5
        y = 2.05 + row * 1.22
        _add_rect(s, x, y, 6.2, 1.08, RGBColor(0x1e, 0x29, 0x3b))
        _add_rect(s, x, y, 0.65, 1.08, DARKGRAY)
        _add_text(s, ref, x + 0.03, y + 0.28, 0.58, 0.48, font_size=11, bold=True, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
        _add_text(s, title, x + 0.75, y + 0.05, 5.3, 0.4, font_size=11, bold=True, color=WHITE)
        _add_text(s, desc,  x + 0.75, y + 0.55, 5.3, 0.42, font_size=9.5, color=GRAY)
    _footer(s, 29)

    # ── SLIDE 30: Q&A / Thank You ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _dark_slide(s)
    _add_rect(s, 0, 0, 13.33, 0.18, GOLD)
    _add_rect(s, 0, 7.32, 13.33, 0.18, RED)
    _add_rect(s, 0, 0, 0.18, 7.5, BLUE)
    _add_text(s, "Thank You", 0.5, 1.3, 12.3, 1.4, font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "We Welcome Your Questions", 0.5, 2.8, 12.3, 0.7, font_size=22, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
    # Summary box
    _add_rect(s, 1.5, 3.7, 10.3, 2.6, RGBColor(0x1e, 0x29, 0x3b))
    _add_text(s, "[ Company Name ]  ·  GSE: [ XXX ]  ·  Recommendation:  [ BUY / SELL / HOLD ]", 1.7, 3.82, 9.9, 0.45, font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    summary_row = [
        ("Current Price", "GHc [ X.XX ]"),
        ("Target Price",  "GHc [ X.XX ]"),
        ("Upside",        "[ +X% ]"),
        ("Horizon",       "12 Months"),
    ]
    for i, (label, val) in enumerate(summary_row):
        x = 1.8 + i * 2.5
        _add_text(s, label, x, 4.4, 2.3, 0.35, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
        _add_text(s, val,   x, 4.78, 2.3, 0.55, font_size=16, bold=True, color=GREEN if i==2 else WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "Built with InvestIQ Professional Calculators  ·  investright.onrender.com", 0.5, 6.45, 12.3, 0.35, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
    _add_text(s, "Young Investors Network  ·  financescholarsyin@gmail.com", 0.5, 6.82, 12.3, 0.35, font_size=10, color=LIGHTBLUE, align=PP_ALIGN.CENTER)
    _footer(s, 30)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


if __name__ == '__main__':
    buf = build_pptx()
    with open('YIN_Stock_Pitch_Template.pptx', 'wb') as f:
        f.write(buf.read())
    print("Done — YIN_Stock_Pitch_Template.pptx created")
